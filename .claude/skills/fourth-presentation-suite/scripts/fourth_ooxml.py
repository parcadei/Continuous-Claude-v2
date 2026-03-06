"""
Fourth Presentation Suite -- OOXML Effects Engine
Low-level DrawingML XML injection for visual effects that python-pptx
cannot produce through its high-level API.

All methods operate on shape._element.spPr (shape properties) or
slide background XML elements. They inject raw OOXML/DrawingML to
produce radial gradients, glow effects, shadows, semi-transparent fills,
rounded corners, and glass-morphism composites.

Designed to run in Claude.ai's sandbox (python-pptx + lxml available).

Usage:
    from fourth_ooxml import OoxmlEffects

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ...)
    OoxmlEffects.apply_glass_card(shape, fill_color='FFFFFF', fill_alpha=10)
    OoxmlEffects.add_glow(shape, radius_pt=12, hex_color='00B69F')
"""

from __future__ import annotations

from typing import List, Optional, Tuple

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu

# ---------------------------------------------------------------------------
# Unit conversion constants
# ---------------------------------------------------------------------------

EMU_PER_PT = 12700
EMU_PER_INCH = 914400
DEG_FACTOR = 60000      # OOXML degrees = real degrees * 60000
PCT_FACTOR = 1000        # OOXML percentage = real % * 1000
ALPHA_FACTOR = 1000      # OOXML alpha = real % * 1000


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert 6-char hex string to RGBColor."""
    h = hex_color.upper().lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# v6 Glow color rotation
# ---------------------------------------------------------------------------

_GLOW_HEX_ROTATION = ['00B69F', 'FAA51A', '9678B6', '6FB4E3']  # teal, sunrise, purple, sky


def get_next_glow_color(counter: int) -> str:
    """Return the next glow hex color in the 4-color rotation.

    Usage: pass a slide counter (0-based) and get a hex color string
    for use with OoxmlEffects.add_glow().
    """
    return _GLOW_HEX_ROTATION[counter % len(_GLOW_HEX_ROTATION)]


# OOXML requires spPr children in a strict order.  If they appear out of
# order PowerPoint will either silently drop effects or show a repair dialog.
_SPPR_CHILD_ORDER = [
    'a:xfrm',
    'a:custGeom',
    'a:prstGeom',
    # -- fills (mutually exclusive, but order matters if somehow both present)
    'a:noFill',
    'a:solidFill',
    'a:gradFill',
    'a:blipFill',
    'a:pattFill',
    'a:grpFill',
    # -- line
    'a:ln',
    # -- effects
    'a:effectLst',
    'a:effectDag',
    # -- 3D scene / shape
    'a:scene3d',
    'a:sp3d',
    # -- extension list
    'a:extLst',
]


# ---------------------------------------------------------------------------
# OoxmlEffects -- static helper class
# ---------------------------------------------------------------------------

class OoxmlEffects:
    """Static methods for injecting DrawingML XML into shape properties.

    Every method takes a python-pptx shape (or slide) and mutates its
    underlying lxml element tree.  No method returns a new element --
    all mutations are in-place.
    """

    # ---- internal helpers ------------------------------------------------

    @staticmethod
    def _get_spPr(shape):
        """Return the <a:spPr> element for a shape, creating if needed."""
        el = shape._element
        spPr = el.find(qn('a:spPr'))
        if spPr is None:
            # Try sp/spPr (normal shape) or cxnSp/spPr
            spPr = el.find(qn('p:spPr'))
        if spPr is None:
            spPr = etree.SubElement(el, qn('a:spPr'))
        return spPr

    @staticmethod
    def _get_or_create_effectLst(spPr):
        """Return <a:effectLst> under spPr, creating if absent."""
        effectLst = spPr.find(qn('a:effectLst'))
        if effectLst is None:
            effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        return effectLst

    @staticmethod
    def _remove_existing_fills(spPr):
        """Remove any existing fill child from spPr."""
        for tag in ('a:noFill', 'a:solidFill', 'a:gradFill',
                     'a:blipFill', 'a:pattFill', 'a:grpFill'):
            existing = spPr.find(qn(tag))
            if existing is not None:
                spPr.remove(existing)

    @staticmethod
    def _reorder_spPr_children(spPr):
        """Re-sort spPr children to match OOXML required order.

        CRITICAL: Must be called after any XML injection into spPr.
        PowerPoint will report file corruption if children are out of order.
        """
        children = list(spPr)
        for child in children:
            spPr.remove(child)

        def sort_key(elem):
            # Strip namespace to get local tag
            tag = etree.QName(elem.tag).localname
            ns = etree.QName(elem.tag).namespace or ''
            # Build prefixed tag for lookup
            prefix = 'a' if 'drawingml' in ns else 'p'
            full_tag = f'{prefix}:{tag}'
            try:
                return _SPPR_CHILD_ORDER.index(full_tag)
            except ValueError:
                return len(_SPPR_CHILD_ORDER)  # unknown tags go last

        children.sort(key=sort_key)
        for child in children:
            spPr.append(child)

    @staticmethod
    def _make_srgbClr(hex_color: str, alpha_pct: Optional[int] = None):
        """Create <a:srgbClr val="RRGGBB"> with optional <a:alpha>."""
        srgb = etree.Element(qn('a:srgbClr'))
        srgb.set('val', hex_color.upper().lstrip('#'))
        if alpha_pct is not None and alpha_pct < 100:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', str(alpha_pct * ALPHA_FACTOR))
        return srgb

    # ---- fill effects ----------------------------------------------------

    @staticmethod
    def set_radial_gradient(shape, stops: List[Tuple[int, str, int]],
                            center: Tuple[int, int] = (50, 50)):
        """Apply a radial gradient fill to a shape.

        Args:
            shape: python-pptx shape object.
            stops: List of (position_pct, hex_color, alpha_pct).
                   e.g., [(0, '002747', 100), (100, '0C4A7D', 80)]
            center: (left_pct, top_pct) for gradient focal point.
        """
        spPr = OoxmlEffects._get_spPr(shape)
        OoxmlEffects._remove_existing_fills(spPr)

        gradFill = etree.SubElement(spPr, qn('a:gradFill'))
        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        for pos_pct, hex_color, alpha_pct in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(pos_pct * PCT_FACTOR))
            gs.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        # Radial path
        path = etree.SubElement(gradFill, qn('a:path'))
        path.set('path', 'circle')
        fillToRect = etree.SubElement(path, qn('a:fillToRect'))
        fillToRect.set('l', str(center[0] * PCT_FACTOR))
        fillToRect.set('t', str(center[1] * PCT_FACTOR))
        fillToRect.set('r', str(center[0] * PCT_FACTOR))
        fillToRect.set('b', str(center[1] * PCT_FACTOR))

        OoxmlEffects._reorder_spPr_children(spPr)

    @staticmethod
    def set_linear_gradient(shape, stops: List[Tuple[int, str, int]],
                            angle_deg: int = 25):
        """Apply a linear gradient fill to a shape.

        Args:
            shape: python-pptx shape object.
            stops: List of (position_pct, hex_color, alpha_pct).
            angle_deg: Gradient angle in degrees.
        """
        spPr = OoxmlEffects._get_spPr(shape)
        OoxmlEffects._remove_existing_fills(spPr)

        gradFill = etree.SubElement(spPr, qn('a:gradFill'))
        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        for pos_pct, hex_color, alpha_pct in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(pos_pct * PCT_FACTOR))
            gs.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', str(angle_deg * DEG_FACTOR))
        lin.set('scaled', '1')

        OoxmlEffects._reorder_spPr_children(spPr)

    @staticmethod
    def set_semi_transparent_fill(shape, hex_color: str, alpha_pct: int):
        """Apply a semi-transparent solid fill.

        Args:
            shape: python-pptx shape object.
            hex_color: 6-char hex (e.g., 'FFFFFF').
            alpha_pct: Opacity 0-100 (0 = fully transparent).
        """
        spPr = OoxmlEffects._get_spPr(shape)
        OoxmlEffects._remove_existing_fills(spPr)

        solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        solidFill.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        OoxmlEffects._reorder_spPr_children(spPr)

    # ---- shape geometry --------------------------------------------------

    @staticmethod
    def set_rounded_corners(shape, radius_pct: int = 8):
        """Set shape geometry to rounded rectangle with specified corner radius.

        Args:
            shape: python-pptx shape object.
            radius_pct: Corner radius as percentage (0-50). 8 is subtle.
        """
        spPr = OoxmlEffects._get_spPr(shape)

        # Remove existing geometry
        for tag in ('a:prstGeom', 'a:custGeom'):
            existing = spPr.find(qn(tag))
            if existing is not None:
                spPr.remove(existing)

        prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'roundRect')
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        # OOXML roundRect adj is in 1/50000ths (50000 = 50% radius)
        gd.set('fmla', f'val {radius_pct * 1000}')

        OoxmlEffects._reorder_spPr_children(spPr)

    # ---- visual effects --------------------------------------------------

    @staticmethod
    def add_glow(shape, radius_pt: int = 10, hex_color: str = '00B69F',
                 alpha_pct: int = 40):
        """Add a glow effect around the shape.

        Args:
            shape: python-pptx shape object.
            radius_pt: Glow spread radius in points.
            hex_color: Glow color (6-char hex).
            alpha_pct: Glow opacity 0-100.
        """
        spPr = OoxmlEffects._get_spPr(shape)
        effectLst = OoxmlEffects._get_or_create_effectLst(spPr)

        # Remove existing glow
        existing = effectLst.find(qn('a:glow'))
        if existing is not None:
            effectLst.remove(existing)

        glow = etree.SubElement(effectLst, qn('a:glow'))
        glow.set('rad', str(radius_pt * EMU_PER_PT))
        glow.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        OoxmlEffects._reorder_spPr_children(spPr)

    @staticmethod
    def add_outer_shadow(shape, blur_pt: int = 6, dist_pt: int = 3,
                         dir_deg: int = 90, hex_color: str = '000000',
                         alpha_pct: int = 50):
        """Add an outer (drop) shadow to the shape.

        Args:
            shape: python-pptx shape object.
            blur_pt: Shadow blur radius in points.
            dist_pt: Shadow offset distance in points.
            dir_deg: Shadow direction in degrees (0=right, 90=down).
            hex_color: Shadow color.
            alpha_pct: Shadow opacity 0-100.
        """
        spPr = OoxmlEffects._get_spPr(shape)
        effectLst = OoxmlEffects._get_or_create_effectLst(spPr)

        existing = effectLst.find(qn('a:outerShdw'))
        if existing is not None:
            effectLst.remove(existing)

        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', str(blur_pt * EMU_PER_PT))
        outerShdw.set('dist', str(dist_pt * EMU_PER_PT))
        outerShdw.set('dir', str(dir_deg * DEG_FACTOR))
        outerShdw.set('algn', 'bl')
        outerShdw.set('rotWithShape', '0')
        outerShdw.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        OoxmlEffects._reorder_spPr_children(spPr)

    @staticmethod
    def add_inner_shadow(shape, blur_pt: int = 5, dist_pt: int = 2,
                         dir_deg: int = 225, hex_color: str = '000000',
                         alpha_pct: int = 50):
        """Add an inner shadow effect to the shape.

        Args:
            shape: python-pptx shape object.
            blur_pt: Inner shadow blur in points.
            dist_pt: Inner shadow offset in points.
            dir_deg: Direction of inner shadow in degrees.
            hex_color: Shadow color.
            alpha_pct: Shadow opacity 0-100.
        """
        spPr = OoxmlEffects._get_spPr(shape)
        effectLst = OoxmlEffects._get_or_create_effectLst(spPr)

        existing = effectLst.find(qn('a:innerShdw'))
        if existing is not None:
            effectLst.remove(existing)

        innerShdw = etree.SubElement(effectLst, qn('a:innerShdw'))
        innerShdw.set('blurRad', str(blur_pt * EMU_PER_PT))
        innerShdw.set('dist', str(dist_pt * EMU_PER_PT))
        innerShdw.set('dir', str(dir_deg * DEG_FACTOR))
        innerShdw.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        OoxmlEffects._reorder_spPr_children(spPr)

    @staticmethod
    def add_soft_edge(shape, radius_pt: int = 3):
        """Add a soft-edge (feathering) effect to the shape.

        Approximates CSS backdrop-filter: blur() by blurring shape edges.

        Args:
            shape: python-pptx shape object.
            radius_pt: Edge blur radius in points.
        """
        spPr = OoxmlEffects._get_spPr(shape)
        effectLst = OoxmlEffects._get_or_create_effectLst(spPr)

        existing = effectLst.find(qn('a:softEdge'))
        if existing is not None:
            effectLst.remove(existing)

        softEdge = etree.SubElement(effectLst, qn('a:softEdge'))
        softEdge.set('rad', str(radius_pt * EMU_PER_PT))

        OoxmlEffects._reorder_spPr_children(spPr)

    # ---- line / border ---------------------------------------------------

    @staticmethod
    def set_line(shape, hex_color: str = 'FFFFFF', alpha_pct: int = 20,
                 width_pt: float = 1.0):
        """Set shape outline with optional transparency.

        Args:
            shape: python-pptx shape object.
            hex_color: Border color.
            alpha_pct: Border opacity 0-100.
            width_pt: Border width in points.
        """
        spPr = OoxmlEffects._get_spPr(shape)

        # Remove existing line
        existing = spPr.find(qn('a:ln'))
        if existing is not None:
            spPr.remove(existing)

        ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w', str(int(width_pt * EMU_PER_PT)))

        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        solidFill.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        OoxmlEffects._reorder_spPr_children(spPr)

    @staticmethod
    def remove_line(shape):
        """Remove shape outline entirely."""
        spPr = OoxmlEffects._get_spPr(shape)
        existing = spPr.find(qn('a:ln'))
        if existing is not None:
            spPr.remove(existing)
        ln = etree.SubElement(spPr, qn('a:ln'))
        etree.SubElement(ln, qn('a:noFill'))
        OoxmlEffects._reorder_spPr_children(spPr)

    # ---- slide background effects ----------------------------------------

    @staticmethod
    def _get_or_create_bg(slide):
        """Get or create the proper <p:bg><p:bgPr> structure on a slide.

        python-pptx's slide.background._element returns <p:cSld>, NOT <p:bg>.
        We must manually create the correct structure:
            <p:cSld>
                <p:bg>           <- wrapper
                    <p:bgPr>     <- properties
                        ...fills...
                    </p:bgPr>
                </p:bg>
                <p:spTree>...</p:spTree>
            </p:cSld>

        The <p:bg> element MUST come before <p:spTree>.
        """
        sld_elem = slide._element
        csld = sld_elem.find(qn('p:cSld'))

        bg = csld.find(qn('p:bg'))
        if bg is None:
            bg = etree.SubElement(csld, qn('p:bg'))
            # Move bg to be the FIRST child of cSld (before spTree)
            csld.remove(bg)
            csld.insert(0, bg)

        bgPr = bg.find(qn('p:bgPr'))
        if bgPr is None:
            bgPr = etree.SubElement(bg, qn('p:bgPr'))

        return bgPr

    @staticmethod
    def set_slide_gradient_bg(slide, stops: List[Tuple[int, str, int]],
                               gradient_type: str = 'radial',
                               center: Tuple[int, int] = (50, 50),
                               angle_deg: int = 25):
        """Apply a gradient background to a slide.

        Args:
            slide: python-pptx slide object.
            stops: List of (position_pct, hex_color, alpha_pct).
            gradient_type: 'radial' or 'linear'.
            center: Focal point for radial (left_pct, top_pct).
            angle_deg: Angle for linear gradients.
        """
        bgPr = OoxmlEffects._get_or_create_bg(slide)

        # Remove existing fills
        for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill'):
            existing = bgPr.find(qn(tag))
            if existing is not None:
                bgPr.remove(existing)

        gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        for pos_pct, hex_color, alpha_pct in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(pos_pct * PCT_FACTOR))
            gs.append(OoxmlEffects._make_srgbClr(hex_color, alpha_pct))

        if gradient_type == 'radial':
            path = etree.SubElement(gradFill, qn('a:path'))
            path.set('path', 'circle')
            fillToRect = etree.SubElement(path, qn('a:fillToRect'))
            fillToRect.set('l', str(center[0] * PCT_FACTOR))
            fillToRect.set('t', str(center[1] * PCT_FACTOR))
            fillToRect.set('r', str(center[0] * PCT_FACTOR))
            fillToRect.set('b', str(center[1] * PCT_FACTOR))
        else:
            lin = etree.SubElement(gradFill, qn('a:lin'))
            lin.set('ang', str(angle_deg * DEG_FACTOR))
            lin.set('scaled', '1')

        if bgPr.find(qn('a:effectLst')) is None:
            etree.SubElement(bgPr, qn('a:effectLst'))

    @staticmethod
    def set_slide_solid_bg(slide, hex_color: str):
        """Apply a solid background color to a slide.

        Args:
            slide: python-pptx slide object.
            hex_color: 6-char hex color.
        """
        bgPr = OoxmlEffects._get_or_create_bg(slide)

        for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill'):
            existing = bgPr.find(qn(tag))
            if existing is not None:
                bgPr.remove(existing)

        solidFill = etree.SubElement(bgPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', hex_color.upper().lstrip('#'))

        if bgPr.find(qn('a:effectLst')) is None:
            etree.SubElement(bgPr, qn('a:effectLst'))

    # ---- composite effects -----------------------------------------------

    @staticmethod
    def apply_glass_card(shape, fill_color: str = 'FFFFFF',
                         fill_alpha: int = 10,
                         border_color: str = 'FFFFFF',
                         border_alpha: int = 15,
                         corner_radius: int = 8,
                         glow_color: Optional[str] = None,
                         glow_radius_pt: int = 12,
                         glow_alpha: int = 35,
                         shadow: bool = True):
        """Apply full glass-morphism composite effect to a shape.

        Combines: rounded corners + semi-transparent fill + alpha border
        + optional glow + outer shadow.

        Args:
            shape: python-pptx shape object.
            fill_color: Card fill color hex.
            fill_alpha: Card fill opacity (0-100).
            border_color: Border color hex.
            border_alpha: Border opacity (0-100).
            corner_radius: Corner radius percentage (0-50).
            glow_color: Optional glow color hex (None = no glow).
            glow_radius_pt: Glow spread in points.
            glow_alpha: Glow opacity (0-100).
            shadow: Whether to add outer shadow.
        """
        # 1. Rounded corners
        OoxmlEffects.set_rounded_corners(shape, corner_radius)

        # 2. Semi-transparent fill
        OoxmlEffects.set_semi_transparent_fill(shape, fill_color, fill_alpha)

        # 3. Alpha border
        OoxmlEffects.set_line(shape, border_color, border_alpha, width_pt=1.0)

        # 4. Optional glow
        if glow_color:
            OoxmlEffects.add_glow(shape, glow_radius_pt, glow_color,
                                  glow_alpha)

        # 5. Outer shadow
        if shadow:
            OoxmlEffects.add_outer_shadow(
                shape,
                blur_pt=8,
                dist_pt=4,
                dir_deg=90,
                hex_color='000000',
                alpha_pct=40,
            )

    @staticmethod
    def apply_status_border(shape, status_color: str, width_pt: float = 2.5):
        """Apply a colored top border to indicate status on a glass card.

        Creates a visible status indicator by setting a thicker, more
        opaque line in the status color.

        Args:
            shape: python-pptx shape object.
            status_color: Hex color for the status border.
            width_pt: Border width in points.
        """
        OoxmlEffects.set_line(shape, status_color, alpha_pct=85,
                              width_pt=width_pt)

    # ------------------------------------------------------------------ #
    #  Chart-level OOXML injection (operates on <c:ser> elements)         #
    # ------------------------------------------------------------------ #

    @staticmethod
    def set_series_gradient_fill(
        ser_element,
        color_hex: str = '00B69F',
        top_alpha_pct: int = 60,
        bottom_alpha_pct: int = 0,
    ):
        """Inject a vertical gradient fill into a chart series element.

        Creates a top-to-bottom gradient from a semi-transparent color
        to fully transparent, producing the teal-to-transparent area fill
        used by Fourth brand area charts.

        Args:
            ser_element: lxml element for the ``<c:ser>`` series node.
            color_hex: Six-character hex color (no ``#`` prefix).
            top_alpha_pct: Opacity at the top of the fill (0-100).
            bottom_alpha_pct: Opacity at the bottom of the fill (0-100).
        """
        # Find or create <c:spPr>
        spPr = ser_element.find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(ser_element, qn('c:spPr'))

        # Remove existing fills
        for fill_tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:pattFill'):
            existing = spPr.find(qn(fill_tag))
            if existing is not None:
                spPr.remove(existing)

        # Build <a:gradFill rotWithShape="0">
        grad_fill = etree.SubElement(spPr, qn('a:gradFill'))
        grad_fill.set('rotWithShape', '0')

        gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))

        # Stop at position 0 (top)
        gs_top = etree.SubElement(gs_lst, qn('a:gs'))
        gs_top.set('pos', '0')
        clr_top = etree.SubElement(gs_top, qn('a:srgbClr'))
        clr_top.set('val', color_hex)
        alpha_top = etree.SubElement(clr_top, qn('a:alpha'))
        alpha_top.set('val', str(top_alpha_pct * ALPHA_FACTOR))

        # Stop at position 100000 (bottom)
        gs_bot = etree.SubElement(gs_lst, qn('a:gs'))
        gs_bot.set('pos', '100000')
        clr_bot = etree.SubElement(gs_bot, qn('a:srgbClr'))
        clr_bot.set('val', color_hex)
        alpha_bot = etree.SubElement(clr_bot, qn('a:alpha'))
        alpha_bot.set('val', str(bottom_alpha_pct * ALPHA_FACTOR))

        # Linear angle: 90 degrees * DEG_FACTOR = 5400000 (top-to-bottom)
        lin = etree.SubElement(grad_fill, qn('a:lin'))
        lin.set('ang', str(90 * DEG_FACTOR))
        lin.set('scaled', '1')

    @staticmethod
    def add_markers_to_series(
        ser_element,
        symbol: str = 'circle',
        size: int = 8,
        fill_hex: str = '00B69F',
        outline_hex: str = 'FFFFFF',
        outline_width_pt: float = 1.5,
    ):
        """Add data-point markers with fill and outline to a chart series.

        Injects a ``<c:marker>`` element with the specified symbol, size,
        fill color, and outline style.

        Args:
            ser_element: lxml element for the ``<c:ser>`` series node.
            symbol: Marker symbol name (e.g. ``'circle'``, ``'square'``).
            size: Marker size (2-72).
            fill_hex: Six-character hex color for the marker fill.
            outline_hex: Six-character hex color for the marker outline.
            outline_width_pt: Outline width in points.
        """
        # Remove existing <c:marker>
        existing = ser_element.find(qn('c:marker'))
        if existing is not None:
            ser_element.remove(existing)

        marker = etree.SubElement(ser_element, qn('c:marker'))

        sym = etree.SubElement(marker, qn('c:symbol'))
        sym.set('val', symbol)

        sz = etree.SubElement(marker, qn('c:size'))
        sz.set('val', str(size))

        spPr = etree.SubElement(marker, qn('c:spPr'))

        # Fill
        solid_fill = etree.SubElement(spPr, qn('a:solidFill'))
        clr = etree.SubElement(solid_fill, qn('a:srgbClr'))
        clr.set('val', fill_hex)

        # Outline
        ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w', str(int(outline_width_pt * EMU_PER_PT)))
        ln_fill = etree.SubElement(ln, qn('a:solidFill'))
        ln_clr = etree.SubElement(ln_fill, qn('a:srgbClr'))
        ln_clr.set('val', outline_hex)

    @staticmethod
    def enable_smooth_lines(ser_element):
        """Enable smooth (spline) interpolation on a chart series line.

        Appends ``<c:smooth val="1"/>`` to the series element, causing
        PowerPoint to render the line with Catmull-Rom spline smoothing.

        Args:
            ser_element: lxml element for the ``<c:ser>`` series node.
        """
        existing = ser_element.find(qn('c:smooth'))
        if existing is not None:
            ser_element.remove(existing)

        smooth = etree.SubElement(ser_element, qn('c:smooth'))
        smooth.set('val', '1')

    @staticmethod
    def set_series_line_style(
        ser_element,
        color_hex: str = '00B69F',
        width_pt: float = 2.5,
    ):
        """Set the visible line style on a chart series.

        Configures the series outline (the drawn line on top of an area
        or line chart) with a solid color and specified width.

        Args:
            ser_element: lxml element for the ``<c:ser>`` series node.
            color_hex: Six-character hex color for the line.
            width_pt: Line width in points.
        """
        # Find or create <c:spPr>
        spPr = ser_element.find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(ser_element, qn('c:spPr'))

        # Remove existing <a:ln>
        existing_ln = spPr.find(qn('a:ln'))
        if existing_ln is not None:
            spPr.remove(existing_ln)

        # Create <a:ln w="...">
        ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w', str(int(width_pt * EMU_PER_PT)))

        solid_fill = etree.SubElement(ln, qn('a:solidFill'))
        clr = etree.SubElement(solid_fill, qn('a:srgbClr'))
        clr.set('val', color_hex)

    # ---- v5 new visual primitives ----------------------------------------

    @staticmethod
    def set_letter_spacing(run, spacing_hundredths: int = 200):
        """Set letter spacing on a text run.

        Args:
            run: A python-pptx Run object.
            spacing_hundredths: Spacing value in hundredths of a point.
                200 = +2pt tracking.  Negative values tighten.
        """
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(spacing_hundredths))

    @staticmethod
    def set_diagonal_gradient(shape, stops, angle_deg=135):
        """Set a diagonal gradient fill on any shape.

        Args:
            shape: python-pptx shape object.
            stops: List of dicts with 'pos' (0-100), 'color' (hex),
                   'alpha' (0-100).
                   e.g., [{'pos': 0, 'color': '00B69F', 'alpha': 100}, ...]
            angle_deg: Gradient angle in degrees. Default 135 (diagonal).
        """
        spPr = OoxmlEffects._get_spPr(shape)
        OoxmlEffects._remove_existing_fills(spPr)

        gradFill = etree.SubElement(spPr, qn('a:gradFill'))
        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        for stop in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(stop['pos'] * PCT_FACTOR))
            gs.append(OoxmlEffects._make_srgbClr(
                stop['color'],
                stop.get('alpha', 100),
            ))

        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', str(angle_deg * DEG_FACTOR))
        lin.set('scaled', '1')

        OoxmlEffects._reorder_spPr_children(spPr)

    @staticmethod
    def create_progress_bar(slide, left, top, width, height,
                            fill_pct, track_color, fill_color,
                            track_alpha=20, corner_radius_pct=50):
        """Create a progress bar from two overlapping shapes.

        Shape 1 (track): Full width, track_color with track_alpha% opacity.
        Shape 2 (fill): Width = fill_pct% of track, fill_color solid.

        Args:
            slide: python-pptx slide object.
            left: Left position (EMU).
            top: Top position (EMU).
            width: Total track width (EMU).
            height: Bar height (EMU).
            fill_pct: Fill percentage 0-100.
            track_color: Hex color for the background track.
            fill_color: Hex color for the filled portion.
            track_alpha: Track opacity 0-100. Default 20.
            corner_radius_pct: Corner radius percentage. Default 50 (full pill).

        Returns:
            (track_shape, fill_shape) tuple.
        """
        from pptx.enum.shapes import MSO_SHAPE

        # Track (full width)
        track = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        OoxmlEffects.set_semi_transparent_fill(track, track_color, track_alpha)
        OoxmlEffects.set_rounded_corners(track, corner_radius_pct)
        OoxmlEffects.remove_line(track)

        # Fill (percentage of width)
        fill_width = int(width * fill_pct / 100)
        fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max(fill_width, 0), height,
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        OoxmlEffects.set_rounded_corners(fill, corner_radius_pct)
        OoxmlEffects.remove_line(fill)

        return (track, fill)

    @staticmethod
    def create_ring_segment(slide, cx, cy, outer_r, inner_r,
                            start_deg, sweep_deg, fill_color, alpha=100):
        """Create a donut/ring arc segment using custom geometry.

        Draws a partial ring (donut arc) at the specified center with
        inner and outer radii.

        Args:
            slide: python-pptx slide object.
            cx: Center X position (EMU).
            cy: Center Y position (EMU).
            outer_r: Outer radius (EMU).
            inner_r: Inner radius (EMU).
            start_deg: Arc start angle in degrees (0 = right, 90 = down).
            sweep_deg: Arc sweep angle in degrees.
            fill_color: Hex color for fill.
            alpha: Fill opacity 0-100. Default 100.

        Returns:
            The created shape.
        """
        from pptx.enum.shapes import MSO_SHAPE

        # Bounding box for the shape
        shape_left = cx - outer_r
        shape_top = cy - outer_r
        shape_size = outer_r * 2

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, shape_left, shape_top, shape_size, shape_size,
        )

        # Replace preset geometry with custom geometry
        spPr = OoxmlEffects._get_spPr(shape)

        # Remove preset geometry
        for tag in ('a:prstGeom', 'a:custGeom'):
            existing = spPr.find(qn(tag))
            if existing is not None:
                spPr.remove(existing)

        # Build custom geometry
        custGeom = etree.SubElement(spPr, qn('a:custGeom'))
        etree.SubElement(custGeom, qn('a:avLst'))
        etree.SubElement(custGeom, qn('a:gdLst'))
        etree.SubElement(custGeom, qn('a:ahLst'))
        cxnLst = etree.SubElement(custGeom, qn('a:cxnLst'))

        # Define coordinate space matching shape size
        rect = etree.SubElement(custGeom, qn('a:rect'))
        rect.set('l', '0')
        rect.set('t', '0')
        rect.set('r', str(shape_size))
        rect.set('b', str(shape_size))

        pathLst = etree.SubElement(custGeom, qn('a:pathLst'))
        path = etree.SubElement(pathLst, qn('a:path'))
        path.set('w', str(shape_size))
        path.set('h', str(shape_size))

        # Calculate arc geometry
        # Center of coordinate space = outer_r, outer_r
        center = outer_r
        import math

        # Start point on outer arc
        start_rad = math.radians(start_deg)
        ox_start = int(center + outer_r * math.cos(start_rad))
        oy_start = int(center + outer_r * math.sin(start_rad))

        # Move to start of outer arc
        moveTo = etree.SubElement(path, qn('a:moveTo'))
        pt = etree.SubElement(moveTo, qn('a:pt'))
        pt.set('x', str(ox_start))
        pt.set('y', str(oy_start))

        # Outer arc
        arcTo_outer = etree.SubElement(path, qn('a:arcTo'))
        arcTo_outer.set('wR', str(outer_r))
        arcTo_outer.set('hR', str(outer_r))
        arcTo_outer.set('stAng', str(int(start_deg * DEG_FACTOR)))
        arcTo_outer.set('swAng', str(int(sweep_deg * DEG_FACTOR)))

        # Line to inner arc end point
        end_deg = start_deg + sweep_deg
        end_rad = math.radians(end_deg)
        ix_end = int(center + inner_r * math.cos(end_rad))
        iy_end = int(center + inner_r * math.sin(end_rad))

        lnTo1 = etree.SubElement(path, qn('a:lnTo'))
        pt1 = etree.SubElement(lnTo1, qn('a:pt'))
        pt1.set('x', str(ix_end))
        pt1.set('y', str(iy_end))

        # Inner arc (reverse direction)
        arcTo_inner = etree.SubElement(path, qn('a:arcTo'))
        arcTo_inner.set('wR', str(inner_r))
        arcTo_inner.set('hR', str(inner_r))
        arcTo_inner.set('stAng', str(int(end_deg * DEG_FACTOR)))
        arcTo_inner.set('swAng', str(int(-sweep_deg * DEG_FACTOR)))

        # Close path
        etree.SubElement(path, qn('a:close'))

        # Apply fill
        OoxmlEffects._remove_existing_fills(spPr)
        solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        solidFill.append(OoxmlEffects._make_srgbClr(fill_color, alpha))

        # Remove line
        existing_ln = spPr.find(qn('a:ln'))
        if existing_ln is not None:
            spPr.remove(existing_ln)
        ln = etree.SubElement(spPr, qn('a:ln'))
        etree.SubElement(ln, qn('a:noFill'))

        OoxmlEffects._reorder_spPr_children(spPr)
        return shape

    @staticmethod
    def add_decorative_line(slide, start_left, start_top, end_left, end_top,
                            color, alpha=10, width_pt=0.5):
        """Add a thin decorative diagonal line shape.

        Used for subtle depth behind content areas.

        Args:
            slide: python-pptx slide object.
            start_left: Line start X position (EMU).
            start_top: Line start Y position (EMU).
            end_left: Line end X position (EMU).
            end_top: Line end Y position (EMU).
            color: Hex color for the line.
            alpha: Line opacity 0-100. Default 10.
            width_pt: Line width in points. Default 0.5.

        Returns:
            The created connector/line shape.
        """
        from pptx.enum.shapes import MSO_SHAPE

        # Calculate bounding box
        left = min(start_left, end_left)
        top = min(start_top, end_top)
        width = abs(end_left - start_left) or Emu(1)
        height = abs(end_top - start_top) or Emu(1)

        # Use a thin rectangle as a line shape (more reliable than connectors
        # for diagonal decorative lines across slides)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )

        # Make the rectangle itself invisible (no fill)
        spPr = OoxmlEffects._get_spPr(shape)
        OoxmlEffects._remove_existing_fills(spPr)
        noFill = etree.SubElement(spPr, qn('a:noFill'))

        # Set the line/border as the visual element
        existing_ln = spPr.find(qn('a:ln'))
        if existing_ln is not None:
            spPr.remove(existing_ln)

        ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w', str(int(width_pt * EMU_PER_PT)))
        ln_fill = etree.SubElement(ln, qn('a:solidFill'))
        ln_fill.append(OoxmlEffects._make_srgbClr(color, alpha))

        OoxmlEffects._reorder_spPr_children(spPr)
        return shape
