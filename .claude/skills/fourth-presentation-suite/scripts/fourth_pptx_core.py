"""
Fourth Presentation Suite -- Core PPTX Generator
Generates Fourth-branded PowerPoint presentations using python-pptx.

Designed to run in Claude.ai's sandbox (python-pptx available, no pip, no network).

Usage:
    from fourth_pptx_core import PresentationBuilder, PresentationRebrander

    builder = PresentationBuilder()
    builder.add_title_slide("Quarterly Review", "Q4 2025 Results")
    builder.add_content_slide(
        "Key Metrics",
        ["Revenue up 15%", "Customer satisfaction at 94%", "New locations: 230"],
    )
    builder.save("quarterly-review.pptx")

Reference: Fourth Brand Essentials v1.0
"""

from __future__ import annotations

import base64
import copy
import itertools
import math
import os
import tempfile
import warnings
from typing import Any, Dict, List, Optional, Tuple, Union

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.util import Emu, Inches, Pt

from fourth_ooxml import OoxmlEffects


# ---------------------------------------------------------------------------
# Class 1: FourthBrand -- brand constants
# ---------------------------------------------------------------------------

class FourthBrand:
    """Fourth visual-identity constants: colours, typography, dimensions."""

    # -- Primary palette -------------------------------------------------- #
    DEEP_BLUE = RGBColor(0x0C, 0x4A, 0x7D)
    TEAL_GREEN = RGBColor(0x00, 0xB6, 0x9F)
    SKY_BLUE = RGBColor(0x6F, 0xB4, 0xE3)

    # -- Secondary palette ------------------------------------------------ #
    MIDNIGHT_NAVY = RGBColor(0x00, 0x27, 0x47)
    DARK_GRAY = RGBColor(0x37, 0x3E, 0x42)
    COOL_GREY = RGBColor(0xCF, 0xD1, 0xD1)
    SOFT_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # -- Tertiary palette (use sparingly) --------------------------------- #
    HOT_RED = RGBColor(0xD8, 0x16, 0x32)
    PURPLE = RGBColor(0x96, 0x78, 0xB6)
    SUNRISE = RGBColor(0xFA, 0xA5, 0x1A)

    # -- v6 Four-color glow rotation (teal -> sunrise -> purple -> sky) --- #
    GLOW_COLORS = [TEAL_GREEN, SUNRISE, PURPLE, SKY_BLUE]

    # -- Chart colour cycle (brand order) --------------------------------- #
    CHART_COLORS = [DEEP_BLUE, TEAL_GREEN, SKY_BLUE, PURPLE, HOT_RED]

    # -- All brand colours for nearest-match rebranding ------------------- #
    ALL_COLORS = [
        DEEP_BLUE, TEAL_GREEN, SKY_BLUE, MIDNIGHT_NAVY,
        DARK_GRAY, COOL_GREY, SOFT_WHITE, WHITE, HOT_RED, PURPLE,
    ]

    # -- Typography ------------------------------------------------------- #
    FONT_FAMILY = "Poppins"
    H1_SIZE = Pt(48)
    H2_SIZE = Pt(36)
    H3_SIZE = Pt(24)
    H4_SIZE = Pt(20)
    BODY_SIZE = Pt(18)
    BODY_SMALL = Pt(14)
    CAPTION_SIZE = Pt(12)
    LINE_SPACING = 1.21  # multiplier

    # -- v2 typography sizes for KPI and section breaks --------------------- #
    KPI_SIZE = Pt(72)          # Full-slide KPI numbers
    KPI_LARGE_SIZE = Pt(96)    # Extra-large KPI variant
    SECTION_TITLE_SIZE = Pt(60)  # Section break headlines (60-80pt)
    CATEGORY_LABEL_SIZE = Pt(12) # Small uppercase category labels

    # Heading weight is Semibold (600) in Poppins.  python-pptx exposes
    # bold as a boolean; we set bold=True to approximate Semibold.
    # Body text uses Regular (400) -- bold=False.

    # -- v3 Surface constants (glass-morphism) ----------------------------- #
    CARD_FILL_COLOR = 'FFFFFF'
    CARD_FILL_ALPHA = 10           # 10% opacity white
    CARD_BORDER_COLOR = 'FFFFFF'
    CARD_BORDER_ALPHA = 15         # 15% opacity border
    GLASS_CORNER_RADIUS = 8        # roundRect adj percentage

    # -- v3 Effect constants ----------------------------------------------- #
    GLOW_RADIUS_PT = 12
    GLOW_ALPHA = 35
    SHADOW_BLUR_PT = 8
    SHADOW_DIST_PT = 4
    SHADOW_ALPHA = 40

    # -- v3 Status colors (hex strings for OOXML injection) --------------- #
    STATUS_GOOD = '00B69F'         # teal
    STATUS_WATCH = 'FFB700'        # amber
    STATUS_ACTION = 'D81632'       # red (corrected per brand-essentials)

    # -- Amber as RGBColor (for python-pptx high-level API) --------------- #
    AMBER = RGBColor(0xFF, 0xB7, 0x00)

    # -- v3 Hex color strings (for OoxmlEffects calls) -------------------- #
    HEX_NAVY = '002747'
    HEX_DEEP_BLUE = '0C4A7D'
    HEX_TEAL = '00B69F'
    HEX_SKY_BLUE = '6FB4E3'
    HEX_WHITE = 'FFFFFF'
    HEX_DARK_GRAY = '373E42'
    HEX_SUNRISE = 'FAA51A'
    HEX_PURPLE = '9678B6'
    HEX_HOT_RED = 'D81632'
    HEX_FOURTH_MIDNIGHT = '0A1929'

    # -- v6 Glow hex rotation (for OoxmlEffects calls) ------------------- #
    HEX_GLOW_COLORS = [HEX_TEAL, HEX_SUNRISE, HEX_PURPLE, HEX_SKY_BLUE]

    # -- Slide dimensions (16:9) ------------------------------------------ #
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # -- Layout margins --------------------------------------------------- #
    MARGIN_LEFT = Inches(0.8)
    MARGIN_TOP = Inches(0.8)
    MARGIN_RIGHT = Inches(0.8)
    MARGIN_BOTTOM = Inches(0.5)
    CONTENT_WIDTH = Inches(11.733)   # SLIDE_WIDTH - left - right
    CONTENT_HEIGHT = Inches(6.2)     # SLIDE_HEIGHT - top - bottom


# ---------------------------------------------------------------------------
# Class 2: TextFormatter -- reusable text-formatting helpers
# ---------------------------------------------------------------------------

class TextFormatter:
    """Static helpers for applying Fourth typography to text frames."""

    # Internal mapping from heading level to (size, bold)
    _LEVEL_MAP = {
        1: (FourthBrand.H1_SIZE, True),
        2: (FourthBrand.H2_SIZE, True),
        3: (FourthBrand.H3_SIZE, True),
        4: (FourthBrand.H4_SIZE, True),
    }

    @staticmethod
    def _apply_line_spacing(paragraph, multiplier: float = FourthBrand.LINE_SPACING):
        """Set line spacing via oxml (python-pptx has limited direct support)."""
        pPr = paragraph._p.get_or_add_pPr()
        lnSpc = pPr.find(qn("a:lnSpc"))
        if lnSpc is None:
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        # Remove existing children
        for child in list(lnSpc):
            lnSpc.remove(child)
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(multiplier * 100000)))

    @staticmethod
    def format_title(
        text_frame,
        text: str,
        level: int = 1,
        color: Optional[RGBColor] = None,
    ):
        """Apply heading formatting.

        Args:
            text_frame: A python-pptx TextFrame.
            text: The heading string.
            level: 1 = H1, 2 = H2, 3 = H3, 4 = H4.
            color: Override colour (default MIDNIGHT_NAVY).
        """
        color = color or FourthBrand.MIDNIGHT_NAVY
        size, bold = TextFormatter._LEVEL_MAP.get(level, (FourthBrand.H2_SIZE, True))

        text_frame.clear()
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = FourthBrand.FONT_FAMILY
        run.font.size = size
        run.font.bold = bold
        run.font.italic = False
        run.font.color.rgb = color
        TextFormatter._apply_line_spacing(p)

    @staticmethod
    def format_body(
        text_frame,
        text: str,
        color: Optional[RGBColor] = None,
        alignment=PP_ALIGN.LEFT,
        size=None,
    ):
        """Apply body-text formatting (Poppins Regular, 1.21x spacing)."""
        color = color or FourthBrand.DARK_GRAY
        font_size = size or FourthBrand.BODY_SIZE

        text_frame.clear()
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0]
        run.font.name = FourthBrand.FONT_FAMILY
        run.font.size = font_size
        run.font.bold = False
        run.font.italic = False
        run.font.color.rgb = color
        TextFormatter._apply_line_spacing(p)

    @staticmethod
    def format_bullets(
        text_frame,
        bullets: List[str],
        level: int = 0,
        color: Optional[RGBColor] = None,
        size=None,
    ):
        """Add bullet points to *text_frame*.

        Each bullet is a separate paragraph with an XML ``buChar`` element
        (the round bullet character) and the requested indent *level*.
        """
        color = color or FourthBrand.DARK_GRAY
        size = size or FourthBrand.BODY_SIZE

        text_frame.clear()
        text_frame.word_wrap = True

        for idx, bullet_text in enumerate(bullets):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet_text
            p.level = level
            p.alignment = PP_ALIGN.LEFT

            # Apply font
            for run in p.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = size
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = color

            # Set bullet character via oxml
            pPr = p._p.get_or_add_pPr()
            # Remove any existing bullet settings
            for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
                existing = pPr.find(qn(tag))
                if existing is not None:
                    pPr.remove(existing)
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "\u2022")  # bullet character

            # Bullet colour matches Teal accent
            buClr = pPr.find(qn("a:buClr"))
            if buClr is None:
                buClr = etree.SubElement(pPr, qn("a:buClr"))
            for child in list(buClr):
                buClr.remove(child)
            srgbClr = etree.SubElement(buClr, qn("a:srgbClr"))
            srgbClr.set("val", "00B69F")  # Teal

            TextFormatter._apply_line_spacing(p)

    @staticmethod
    def enforce_4x6(bullets: List[str]) -> Tuple[bool, List[str]]:
        """Validate against the v2 content rule: max 4 bullets, max 8 words each.

        Returns:
            (is_valid, list_of_warnings)

        v2 hard limits:
          - Max 4 bullets per slide
          - Max 8 words per bullet
        """
        warns: List[str] = []
        count = len(bullets)
        word_counts = [len(b.split()) for b in bullets]

        # Hard limit: max 4 bullets
        if count > 4:
            warns.append(
                f"Too many bullets ({count}). Max is 4 per slide -- split across slides."
            )

        # Hard limit: max 8 words per bullet
        for i, wc in enumerate(word_counts):
            if wc > 8:
                warns.append(
                    f"Bullet {i + 1} has {wc} words (max 8). Rewrite shorter."
                )

        return (len(warns) == 0, warns)


# ---------------------------------------------------------------------------
# Class 3: BackgroundManager -- slide background helpers
# ---------------------------------------------------------------------------

class BackgroundManager:
    """Static helpers for applying backgrounds to slides."""

    @staticmethod
    def apply_solid(slide, color: RGBColor):
        """Fill the slide background with a solid colour."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    @staticmethod
    def _get_or_create_bgPr(slide):
        """Get or create proper <p:bg><p:bgPr> under <p:cSld>.

        slide.background._element returns <p:cSld>, NOT <p:bg>.
        We must build: <p:cSld><p:bg><p:bgPr>...</p:bgPr></p:bg><p:spTree>...
        """
        sld_elem = slide._element
        csld = sld_elem.find(qn("p:cSld"))
        bg = csld.find(qn("p:bg"))
        if bg is None:
            bg = etree.SubElement(csld, qn("p:bg"))
            csld.remove(bg)
            csld.insert(0, bg)
        bgPr = bg.find(qn("p:bgPr"))
        if bgPr is None:
            bgPr = etree.SubElement(bg, qn("p:bgPr"))
        return bgPr

    @staticmethod
    def apply_gradient(
        slide,
        color1: RGBColor,
        color2: RGBColor,
        angle: int = 0,
    ):
        """Apply a two-stop linear gradient via direct XML manipulation.

        python-pptx's public API does not support gradient backgrounds,
        so we build the ``<a:gradFill>`` element on the slide's background.
        """
        bgPr = BackgroundManager._get_or_create_bgPr(slide)

        # Remove any existing fill
        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill"):
            existing = bgPr.find(qn(tag))
            if existing is not None:
                bgPr.remove(existing)

        # Build gradFill
        gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
        gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

        # Stop 1 at position 0%
        gs0 = etree.SubElement(gsLst, qn("a:gs"))
        gs0.set("pos", "0")
        srgb0 = etree.SubElement(gs0, qn("a:srgbClr"))
        srgb0.set("val", str(color1))

        # Stop 2 at position 100%
        gs1 = etree.SubElement(gsLst, qn("a:gs"))
        gs1.set("pos", "100000")
        srgb1 = etree.SubElement(gs1, qn("a:srgbClr"))
        srgb1.set("val", str(color2))

        # Linear direction
        lin = etree.SubElement(gradFill, qn("a:lin"))
        # OOXML angle is in 60,000ths of a degree
        lin.set("ang", str(angle * 60000))
        lin.set("scaled", "1")

        # Ensure <a:effectLst/> sibling (required by some renderers)
        if bgPr.find(qn("a:effectLst")) is None:
            etree.SubElement(bgPr, qn("a:effectLst"))

    @staticmethod
    def _apply_three_stop_gradient(
        slide,
        stops: List[Tuple[int, RGBColor]],
        angle: int = 0,
    ):
        """Apply an N-stop gradient.  *stops* is a list of (position_pct, colour)."""
        bgPr = BackgroundManager._get_or_create_bgPr(slide)

        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill"):
            existing = bgPr.find(qn(tag))
            if existing is not None:
                bgPr.remove(existing)

        gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
        gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

        for pos_pct, color in stops:
            gs = etree.SubElement(gsLst, qn("a:gs"))
            gs.set("pos", str(pos_pct * 1000))  # position in 1/1000ths of %
            srgb = etree.SubElement(gs, qn("a:srgbClr"))
            srgb.set("val", str(color))

        lin = etree.SubElement(gradFill, qn("a:lin"))
        lin.set("ang", str(angle * 60000))
        lin.set("scaled", "1")

        if bgPr.find(qn("a:effectLst")) is None:
            etree.SubElement(bgPr, qn("a:effectLst"))

    @staticmethod
    def apply_vignette(slide):
        """Apply the signature Fourth vignette gradient.

        Teal (#00B69F) at 0% -> Sky Blue (#6FB4E3) at 40% -> Deep Blue (#0C4A7D) at 100%
        Angle: 25 degrees.
        """
        BackgroundManager._apply_three_stop_gradient(
            slide,
            stops=[
                (0, FourthBrand.TEAL_GREEN),
                (40, FourthBrand.SKY_BLUE),
                (100, FourthBrand.DEEP_BLUE),
            ],
            angle=25,
        )

    @staticmethod
    def apply_image_background(slide, image_path: str):
        """Set an image as the slide background via XML.

        The image is embedded in the presentation and referenced by the
        slide's background ``blipFill``.
        """
        if not os.path.isfile(image_path):
            warnings.warn(f"Background image not found: {image_path}")
            return

        # Add image as a relationship on the slide part
        slide_part = slide.part
        with open(image_path, "rb") as f:
            image_blob = f.read()

        # Determine content type from extension
        ext = os.path.splitext(image_path)[1].lower()
        ct_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".bmp": "image/bmp",
        }
        content_type = ct_map.get(ext, "image/png")

        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        image_part, rId = slide_part.get_or_add_image_part(image_path)

        # Build background XML
        bgPr = BackgroundManager._get_or_create_bgPr(slide)

        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill"):
            existing = bgPr.find(qn(tag))
            if existing is not None:
                bgPr.remove(existing)

        blipFill = etree.SubElement(bgPr, qn("a:blipFill"))
        blipFill.set("dpi", "0")
        blipFill.set("rotWithShape", "1")
        blip = etree.SubElement(blipFill, qn("a:blip"))
        blip.set(qn("r:embed"), rId)
        stretch = etree.SubElement(blipFill, qn("a:stretch"))
        etree.SubElement(stretch, qn("a:fillRect"))

        if bgPr.find(qn("a:effectLst")) is None:
            etree.SubElement(bgPr, qn("a:effectLst"))


# ---------------------------------------------------------------------------
# Class 3b: BackgroundRotation -- enforce background variety
# ---------------------------------------------------------------------------

class BackgroundRotation:
    """Track and enforce background variety across slides.

    Prevents more than 2 consecutive slides from sharing the same
    background treatment category.
    """

    # Background categories
    VIGNETTE = "vignette"
    DARK = "dark"
    LIGHT = "light"
    GRADIENT = "gradient"
    WHITE = "white"

    def __init__(self):
        self._history: List[str] = []

    def next_background(self, slide_type: str) -> str:
        """Return recommended background category for *slide_type*.

        If the recommendation would create 3 consecutive same-category slides,
        it suggests an alternative.
        """
        # Default mapping from slide type to preferred bg
        _type_to_bg = {
            "title": self.VIGNETTE,
            "section": self.GRADIENT,
            "kpi": self.DARK,
            "kpi_grid": self.DARK,
            "content": self.LIGHT,
            "chart": self.WHITE,
            "quote": self.DARK,
            "problem": self.DARK,
            "comparison": self.LIGHT,
            "table": self.WHITE,
            "image": self.LIGHT,
            "closing": self.VIGNETTE,
        }
        preferred = _type_to_bg.get(slide_type, self.LIGHT)

        # Check for 3-consecutive violation
        if (len(self._history) >= 2
                and self._history[-1] == preferred
                and self._history[-2] == preferred):
            # Suggest an alternative
            alternatives = [self.LIGHT, self.DARK, self.WHITE, self.GRADIENT]
            for alt in alternatives:
                if alt != preferred:
                    return alt

        return preferred

    def record(self, bg_category: str):
        """Record the background category that was actually used."""
        self._history.append(bg_category)

    def would_violate(self, bg_category: str) -> bool:
        """Return True if using *bg_category* would create 3 consecutive same."""
        if len(self._history) < 2:
            return False
        return (self._history[-1] == bg_category
                and self._history[-2] == bg_category)

    @property
    def history(self) -> List[str]:
        return list(self._history)


# ---------------------------------------------------------------------------
# Class 4: ImageHandler -- image insertion helpers
# ---------------------------------------------------------------------------

class ImageHandler:
    """Static helpers for inserting images into slides."""

    @staticmethod
    def insert_from_file(
        slide,
        image_path: str,
        left,
        top,
        width,
        height=None,
    ):
        """Insert an image from a file path.

        If *height* is ``None`` the aspect ratio is preserved.
        Returns the Picture shape or ``None`` on failure.
        """
        if not os.path.isfile(image_path):
            warnings.warn(f"Image file not found: {image_path}")
            return None
        try:
            pic = slide.shapes.add_picture(image_path, left, top, width, height)
            return pic
        except Exception as exc:
            warnings.warn(f"Failed to insert image {image_path}: {exc}")
            return None

    @staticmethod
    def insert_from_base64(
        slide,
        b64_data: str,
        left,
        top,
        width,
        height=None,
        ext: str = ".png",
    ):
        """Insert an image from a base64-encoded string.

        Decodes to a temporary file, inserts, then cleans up.
        """
        try:
            raw = base64.b64decode(b64_data)
        except Exception as exc:
            warnings.warn(f"Base64 decode failed: {exc}")
            return None

        tmp_path = None
        try:
            fd, tmp_path = tempfile.mkstemp(suffix=ext)
            os.write(fd, raw)
            os.close(fd)
            pic = slide.shapes.add_picture(tmp_path, left, top, width, height)
            return pic
        except Exception as exc:
            warnings.warn(f"Failed to insert base64 image: {exc}")
            return None
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    @staticmethod
    def create_placeholder(
        slide,
        left,
        top,
        width,
        height,
        label: str = "Image",
    ):
        """Create a branded placeholder rectangle with a centred label."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = FourthBrand.COOL_GREY
        shape.line.color.rgb = FourthBrand.SKY_BLUE
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[ {label} ]"
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.BODY_SMALL
            run.font.color.rgb = FourthBrand.DARK_GRAY
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return shape

    @staticmethod
    def insert_logo(
        slide,
        logo_path: str,
        position: str = "top-right",
        variant: str = "standard",
    ):
        """Insert the Fourth logo at a predefined position.

        Positions: top-right, top-left, bottom-right, bottom-left, bottom-center.
        """
        if not os.path.isfile(logo_path):
            warnings.warn(f"Logo file not found: {logo_path}")
            return None

        logo_w = Inches(1.8)
        pad = Inches(0.4)
        sw = FourthBrand.SLIDE_WIDTH
        sh = FourthBrand.SLIDE_HEIGHT

        positions = {
            "top-right":     (sw - logo_w - pad, pad),
            "top-left":      (pad, pad),
            "bottom-right":  (sw - logo_w - pad, sh - Inches(0.8)),
            "bottom-left":   (pad, sh - Inches(0.8)),
            "bottom-center": ((sw - logo_w) // 2, sh - Inches(0.8)),
        }
        left, top = positions.get(position, positions["top-right"])

        try:
            pic = slide.shapes.add_picture(logo_path, left, top, logo_w)
            return pic
        except Exception as exc:
            warnings.warn(f"Failed to insert logo: {exc}")
            return None


# ---------------------------------------------------------------------------
# Class 5: PresentationBuilder -- main builder API
# ---------------------------------------------------------------------------

class PresentationBuilder:
    """High-level builder for Fourth-branded presentations.

    Every ``add_*`` method returns the created ``Slide`` object so callers
    can layer additional shapes on top.
    """

    def __init__(self, template_path: Optional[str] = None):
        """Create a new presentation (or clone *template_path*)."""
        if template_path and os.path.isfile(template_path):
            self._prs = Presentation(template_path)
        else:
            self._prs = Presentation()

        # Force 16:9 dimensions
        self._prs.slide_width = FourthBrand.SLIDE_WIDTH
        self._prs.slide_height = FourthBrand.SLIDE_HEIGHT

        # Dark background auto-rotation cycle
        self._dark_bg_cycle = itertools.cycle([
            'radial', 'corner_tl', 'corner_br', 'sweep',
        ])

        # Section break variant counter (cycles through 3 variants)
        self._section_break_counter = 0

    # -- internal helpers ------------------------------------------------- #

    def _blank_layout(self):
        """Return the blank slide layout (index 6 in a default template)."""
        layouts = self._prs.slide_layouts
        # Prefer a truly blank layout; fall back to last one
        for idx in (6, len(layouts) - 1, 0):
            try:
                return layouts[idx]
            except IndexError:
                continue
        return layouts[0]

    def _add_slide(self):
        """Add a new blank slide and return it."""
        return self._prs.slides.add_slide(self._blank_layout())

    def _add_textbox(self, slide, left, top, width, height):
        """Shortcut: add a text box and return its text_frame."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        return txBox.text_frame

    @staticmethod
    def _centered_left(width):
        """Left coordinate to horizontally centre an element of *width*."""
        return (FourthBrand.SLIDE_WIDTH - width) // 2

    # -- v3 OOXML-enhanced helpers ---------------------------------------- #

    def _set_dark_bg(self, slide, variant: str = 'navy'):
        """Apply a dark background using OOXML effects.

        Variants:
            'navy': Solid Midnight Navy.
            'radial': Radial gradient from center (50,35).
            'corner_tl': Radial gradient from top-left (15,20).
            'corner_br': Radial gradient from bottom-right (85,80).
            'sweep': Linear gradient at 135 degrees (diagonal dark sweep).
            'auto': Auto-rotate through radial/corner_tl/corner_br/sweep.
            'vignette': Linear gradient 25deg (Teal -> Sky -> Deep Blue).
                        Used for title and closing slides.
        """
        # Auto-rotation: advance the cycle to get varied dark backgrounds
        if variant == 'auto':
            variant = next(self._dark_bg_cycle)

        _RADIAL_STOPS = [
            (0, FourthBrand.HEX_DEEP_BLUE, 100),
            (50, FourthBrand.HEX_NAVY, 100),
            (100, '001020', 100),
        ]

        if variant == 'radial':
            OoxmlEffects.set_slide_gradient_bg(
                slide, stops=_RADIAL_STOPS,
                gradient_type='radial', center=(50, 35),
            )
        elif variant == 'corner_tl':
            OoxmlEffects.set_slide_gradient_bg(
                slide, stops=_RADIAL_STOPS,
                gradient_type='radial', center=(15, 20),
            )
        elif variant == 'corner_br':
            OoxmlEffects.set_slide_gradient_bg(
                slide, stops=_RADIAL_STOPS,
                gradient_type='radial', center=(85, 80),
            )
        elif variant == 'sweep':
            OoxmlEffects.set_slide_gradient_bg(
                slide,
                stops=[
                    (0, FourthBrand.HEX_NAVY, 100),
                    (60, '001020', 100),
                    (100, FourthBrand.HEX_DEEP_BLUE, 100),
                ],
                gradient_type='linear', angle_deg=135,
            )
        elif variant == 'vignette':
            OoxmlEffects.set_slide_gradient_bg(
                slide,
                stops=[
                    (0, FourthBrand.HEX_TEAL, 100),
                    (40, FourthBrand.HEX_SKY_BLUE, 100),
                    (100, FourthBrand.HEX_DEEP_BLUE, 100),
                ],
                gradient_type='linear', angle_deg=25,
            )
        else:
            OoxmlEffects.set_slide_solid_bg(slide, FourthBrand.HEX_NAVY)

    def _add_glass_card(self, slide, left, top, width, height,
                        status_color: Optional[str] = None,
                        glow: bool = False):
        """Add a glass-morphism card shape to the slide.

        Returns the created shape for further manipulation.

        Args:
            slide: The target slide.
            left, top, width, height: Position and size (EMU values).
            status_color: Optional hex color for status border (e.g., '00B69F').
            glow: Whether to add a glow effect in the status color.
        """
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        OoxmlEffects.apply_glass_card(
            card,
            fill_color=FourthBrand.CARD_FILL_COLOR,
            fill_alpha=FourthBrand.CARD_FILL_ALPHA,
            border_color=FourthBrand.CARD_BORDER_COLOR,
            border_alpha=FourthBrand.CARD_BORDER_ALPHA,
            corner_radius=FourthBrand.GLASS_CORNER_RADIUS,
            glow_color=status_color if (glow or status_color) else None,
            glow_radius_pt=FourthBrand.GLOW_RADIUS_PT,
            glow_alpha=FourthBrand.GLOW_ALPHA,
            shadow=True,
        )
        # Status indicated via glow halo only -- no boxy perimeter border
        return card

    # -- public slide methods --------------------------------------------- #

    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        logo_path: Optional[str] = None,
    ):
        """Title slide: vignette background, centred white text, optional logo."""
        slide = self._add_slide()
        self._set_dark_bg(slide, 'vignette')

        # Title -- centred, upper-third
        title_w = FourthBrand.CONTENT_WIDTH
        title_h = Inches(1.8)
        title_top = Inches(2.2)
        tf = self._add_textbox(
            slide,
            self._centered_left(title_w),
            title_top,
            title_w,
            title_h,
        )
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H1_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        # Teal accent line below title
        accent_w = Inches(3.0)
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._centered_left(accent_w),
            title_top + title_h + Inches(0.1),
            accent_w, Inches(0.04),
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        accent_line.line.fill.background()

        # Subtitle
        if subtitle:
            sub_h = Inches(1.0)
            sub_top = title_top + title_h + Inches(0.25)
            tf_sub = self._add_textbox(
                slide,
                self._centered_left(title_w),
                sub_top,
                title_w,
                sub_h,
            )
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.alignment = PP_ALIGN.CENTER
            for run in p_sub.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.H3_SIZE
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE
                OoxmlEffects.set_letter_spacing(run, 100)
            TextFormatter._apply_line_spacing(p_sub)

        # Logo
        if logo_path:
            ImageHandler.insert_logo(slide, logo_path, position="top-right")

        return slide

    def add_content_slide(
        self,
        title: str,
        bullets: Optional[List[str]] = None,
        body_text: Optional[str] = None,
        layout: str = "single",
        category_label: Optional[str] = None,
        bg: str = 'light',
        numbered: bool = False,
    ):
        """Content slide with title and body text or bullets.

        *layout* can be ``"single"`` (standard) or ``"wide"`` (full-width body).
        *category_label*: optional small uppercase teal label above the title
        (e.g., "WORKFORCE INTELLIGENCE", "ROI").
        *bg*: ``'light'`` (white background) or ``'dark'`` (navy with teal glow).
        *numbered*: if True, format bullets as ``01``, ``02`` etc. with teal
        bold numbers instead of dot bullets. Best with ``bg='dark'``.
        """
        slide = self._add_slide()
        is_dark = bg == 'dark'

        if is_dark:
            self._set_dark_bg(slide, 'auto')
        else:
            BackgroundManager.apply_solid(slide, FourthBrand.WHITE)

        # Colors based on background
        title_color = FourthBrand.WHITE if is_dark else FourthBrand.DEEP_BLUE
        body_color = FourthBrand.COOL_GREY if is_dark else None
        label_color = FourthBrand.TEAL_GREEN if is_dark else None
        accent_color = FourthBrand.TEAL_GREEN

        # Category label (above title)
        label_offset = Inches(0)
        if category_label:
            self.add_category_label(slide, category_label, color=label_color)
            label_offset = Inches(0.35)

        # Title bar
        title_h = Inches(1.0)
        title_top = FourthBrand.MARGIN_TOP + label_offset
        tf_title = self._add_textbox(
            slide,
            FourthBrand.MARGIN_LEFT,
            title_top,
            FourthBrand.CONTENT_WIDTH,
            title_h,
        )
        TextFormatter.format_title(tf_title, title, level=2, color=title_color)

        # Teal accent line below title
        line_top = title_top + title_h + Inches(0.05)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            FourthBrand.MARGIN_LEFT,
            line_top,
            Inches(2.0),
            Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_color
        accent.line.fill.background()  # no border
        if is_dark:
            OoxmlEffects.add_glow(accent, radius_pt=6,
                                  hex_color=FourthBrand.HEX_TEAL, alpha_pct=20)

        # Body area
        body_top = line_top + Inches(0.35)
        body_w = FourthBrand.CONTENT_WIDTH if layout == "wide" else Inches(10.5)
        body_h = FourthBrand.SLIDE_HEIGHT - body_top - FourthBrand.MARGIN_BOTTOM

        if bullets and numbered:
            # Numbered items: teal bold "01" prefix + white/dark text
            for idx, b in enumerate(bullets):
                item_top = body_top + Inches(idx * 0.55)
                tf_item = self._add_textbox(
                    slide, FourthBrand.MARGIN_LEFT, item_top,
                    body_w, Inches(0.50),
                )
                tf_item.word_wrap = True
                p = tf_item.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                # Number run
                run_num = p.add_run()
                run_num.text = f"{idx + 1:02d}  "
                run_num.font.name = FourthBrand.FONT_FAMILY
                run_num.font.size = FourthBrand.BODY_SIZE
                run_num.font.bold = True
                run_num.font.color.rgb = FourthBrand.TEAL_GREEN

                # Text run
                run_txt = p.add_run()
                run_txt.text = b
                run_txt.font.name = FourthBrand.FONT_FAMILY
                run_txt.font.size = FourthBrand.BODY_SIZE
                run_txt.font.bold = False
                run_txt.font.color.rgb = (
                    FourthBrand.WHITE if is_dark else FourthBrand.DARK_GRAY
                )
        elif bullets:
            tf_body = self._add_textbox(
                slide, FourthBrand.MARGIN_LEFT, body_top, body_w, body_h,
            )
            # v2: hard-enforce max 4 bullets, max 8 words each
            if len(bullets) > 4:
                warnings.warn(
                    f"4x8 rule: {len(bullets)} bullets (max 4). Truncating."
                )
                bullets = bullets[:4]
            for i, b in enumerate(bullets):
                wc = len(b.split())
                if wc > 8:
                    warnings.warn(
                        f"4x8 rule: bullet {i+1} has {wc} words (max 8)."
                    )
            TextFormatter.format_bullets(tf_body, bullets, color=body_color)
        elif body_text:
            tf_body = self._add_textbox(
                slide, FourthBrand.MARGIN_LEFT, body_top, body_w, body_h,
            )
            TextFormatter.format_body(tf_body, body_text, color=body_color)

        return slide

    def add_section_break(
        self,
        title: str,
        subtitle: Optional[str] = None,
    ):
        """Section divider with 3 rotating visual variants.

        Variant 0: Left-aligned title with decorative diagonal lines on right.
        Variant 1: Full-bleed section number background with centred title.
        Variant 2: Asymmetric layout with gradient shape on right side.

        All variants share a radial gradient background.
        """
        slide = self._add_slide()
        OoxmlEffects.set_slide_gradient_bg(
            slide,
            stops=[
                (0, FourthBrand.HEX_TEAL, 100),
                (50, FourthBrand.HEX_DEEP_BLUE, 100),
                (100, FourthBrand.HEX_NAVY, 100),
            ],
            gradient_type='radial',
            center=(50, 50),
        )

        variant = self._section_break_counter % 3
        self._section_break_counter += 1

        if variant == 0:
            self._section_break_v0(slide, title, subtitle)
        elif variant == 1:
            self._section_break_v1(slide, title, subtitle)
        else:
            self._section_break_v2(slide, title, subtitle)

        return slide

    # -- Section break variant helpers --------------------------------------- #

    def _section_break_v0(self, slide, title, subtitle):
        """Variant 0: Left-aligned title with decorative lines on right."""
        title_w = Inches(7.0)
        title_h = Inches(2.0)
        title_top = Inches(2.2) if not subtitle else Inches(1.8)
        title_left = Inches(1.0)

        tf = self._add_textbox(slide, title_left, title_top, title_w, title_h)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.SECTION_TITLE_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        if subtitle:
            sub_h = Inches(0.8)
            sub_top = title_top + title_h + Inches(0.1)
            tf_sub = self._add_textbox(
                slide, title_left, sub_top, title_w, sub_h,
            )
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.alignment = PP_ALIGN.LEFT
            for run in p_sub.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.H3_SIZE
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE
            TextFormatter._apply_line_spacing(p_sub)

        # Decorative diagonal lines on the right side
        OoxmlEffects.add_decorative_line(
            slide,
            start_left=Inches(8.5), start_top=Inches(1.0),
            end_left=Inches(12.5), end_top=Inches(5.5),
            color='FFFFFF', alpha=10, width_pt=0.75,
        )
        OoxmlEffects.add_decorative_line(
            slide,
            start_left=Inches(9.0), start_top=Inches(0.5),
            end_left=Inches(13.0), end_top=Inches(6.0),
            color='00B69F', alpha=8, width_pt=0.5,
        )
        OoxmlEffects.add_decorative_line(
            slide,
            start_left=Inches(8.0), start_top=Inches(2.0),
            end_left=Inches(12.0), end_top=Inches(7.0),
            color='FFFFFF', alpha=6, width_pt=0.5,
        )

    def _section_break_v1(self, slide, title, subtitle):
        """Variant 1: Full-bleed section number background with centred title."""
        # Section number based on counter (already incremented)
        section_num = self._section_break_counter
        num_text = f"{section_num:02d}"

        # Giant background number (200pt, 5% opacity)
        num_w = Inches(6.0)
        num_h = Inches(4.0)
        num_left = self._centered_left(num_w)
        num_top = Inches(1.5)

        tf_num = self._add_textbox(slide, num_left, num_top, num_w, num_h)
        tf_num.word_wrap = False
        p_num = tf_num.paragraphs[0]
        p_num.text = num_text
        p_num.alignment = PP_ALIGN.CENTER
        for run in p_num.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(200)
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
            # Set text opacity to 5% via OOXML alpha element
            rPr = run._r.get_or_add_rPr()
            solidFill = rPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                    alpha_el.set('val', '5000')  # 5%

        # Title overlaid on top (fully opaque, centred)
        title_w = FourthBrand.CONTENT_WIDTH
        title_h = Inches(2.0)
        title_top = Inches(2.2) if not subtitle else Inches(1.8)
        tf = self._add_textbox(
            slide, self._centered_left(title_w), title_top, title_w, title_h,
        )
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.SECTION_TITLE_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        if subtitle:
            sub_h = Inches(0.8)
            sub_top = title_top + title_h + Inches(0.1)
            tf_sub = self._add_textbox(
                slide, self._centered_left(title_w), sub_top, title_w, sub_h,
            )
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.alignment = PP_ALIGN.CENTER
            for run in p_sub.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.H3_SIZE
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE
            TextFormatter._apply_line_spacing(p_sub)

    def _section_break_v2(self, slide, title, subtitle):
        """Variant 2: Asymmetric layout with gradient shape on right."""
        # Decorative gradient rectangle on the right 40%
        grad_left = Inches(8.0)
        grad_top = Inches(0)
        grad_w = Inches(5.33)
        grad_h = FourthBrand.SLIDE_HEIGHT

        grad_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, grad_left, grad_top, grad_w, grad_h,
        )
        OoxmlEffects.set_diagonal_gradient(
            grad_shape,
            stops=[
                {'pos': 0, 'color': '00B69F', 'alpha': 15},
                {'pos': 100, 'color': '0C4A7D', 'alpha': 5},
            ],
        )
        OoxmlEffects.remove_line(grad_shape)

        # Title: left-aligned, 60% width
        title_w = Inches(7.0)
        title_h = Inches(2.0)
        title_top = Inches(2.5)
        title_left = Inches(1.0)

        tf = self._add_textbox(slide, title_left, title_top, title_w, title_h)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.SECTION_TITLE_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        if subtitle:
            sub_h = Inches(0.8)
            sub_top = title_top + title_h + Inches(0.1)
            tf_sub = self._add_textbox(
                slide, title_left, sub_top, title_w, sub_h,
            )
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.alignment = PP_ALIGN.LEFT
            for run in p_sub.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.H3_SIZE
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE
            TextFormatter._apply_line_spacing(p_sub)

    def add_two_column(
        self,
        title: str,
        left_content: Union[str, List[str]],
        right_content: Union[str, List[str]],
    ):
        """Two-column layout.  Each content arg can be a string or bullet list."""
        slide = self._add_slide()
        BackgroundManager.apply_solid(slide, FourthBrand.WHITE)

        # Title
        title_h = Inches(1.0)
        tf_title = self._add_textbox(
            slide,
            FourthBrand.MARGIN_LEFT,
            FourthBrand.MARGIN_TOP,
            FourthBrand.CONTENT_WIDTH,
            title_h,
        )
        TextFormatter.format_title(tf_title, title, level=2, color=FourthBrand.DEEP_BLUE)

        # Accent line
        line_top = FourthBrand.MARGIN_TOP + title_h + Inches(0.05)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            FourthBrand.MARGIN_LEFT,
            line_top,
            Inches(2.0),
            Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        accent.line.fill.background()

        # Columns
        col_top = line_top + Inches(0.35)
        col_w = Inches(5.4)
        col_h = FourthBrand.SLIDE_HEIGHT - col_top - FourthBrand.MARGIN_BOTTOM
        gutter = Inches(0.93)

        left_left = FourthBrand.MARGIN_LEFT
        right_left = left_left + col_w + gutter

        for content, col_left in [(left_content, left_left), (right_content, right_left)]:
            tf = self._add_textbox(slide, col_left, col_top, col_w, col_h)
            if isinstance(content, list):
                TextFormatter.format_bullets(tf, content)
            else:
                TextFormatter.format_body(tf, content)

        return slide

    def add_three_column(
        self,
        title: str,
        columns: List[Dict[str, str]],
    ):
        """Three-column layout.

        *columns* is a list of up to three dicts, each with ``heading`` and ``body``.
        """
        slide = self._add_slide()
        BackgroundManager.apply_solid(slide, FourthBrand.WHITE)

        # Title
        title_h = Inches(1.0)
        tf_title = self._add_textbox(
            slide,
            FourthBrand.MARGIN_LEFT,
            FourthBrand.MARGIN_TOP,
            FourthBrand.CONTENT_WIDTH,
            title_h,
        )
        TextFormatter.format_title(tf_title, title, level=2, color=FourthBrand.DEEP_BLUE)

        # Accent line
        line_top = FourthBrand.MARGIN_TOP + title_h + Inches(0.05)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            FourthBrand.MARGIN_LEFT,
            line_top,
            Inches(2.0),
            Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        accent.line.fill.background()

        # Column geometry
        col_top = line_top + Inches(0.35)
        col_count = min(len(columns), 3)
        gutter = Inches(0.5)
        total_gutter = gutter * (col_count - 1)
        col_w = (FourthBrand.CONTENT_WIDTH - total_gutter) // col_count
        col_h = FourthBrand.SLIDE_HEIGHT - col_top - FourthBrand.MARGIN_BOTTOM

        for i, col_data in enumerate(columns[:3]):
            col_left = FourthBrand.MARGIN_LEFT + i * (col_w + gutter)

            # Column heading
            head_h = Inches(0.7)
            tf_head = self._add_textbox(slide, col_left, col_top, col_w, head_h)
            TextFormatter.format_title(
                tf_head,
                col_data.get("heading", ""),
                level=4,
                color=FourthBrand.DEEP_BLUE,
            )

            # Column body
            body_top = col_top + head_h + Inches(0.1)
            body_h = col_h - head_h - Inches(0.1)
            tf_body = self._add_textbox(slide, col_left, body_top, col_w, body_h)
            body_text = col_data.get("body", "")
            if isinstance(body_text, list):
                TextFormatter.format_bullets(tf_body, body_text)
            else:
                TextFormatter.format_body(tf_body, body_text)

        return slide

    def add_image_slide(
        self,
        title: str,
        image_path: Optional[str] = None,
        image_b64: Optional[str] = None,
        caption: Optional[str] = None,
    ):
        """Image-focused slide: title at top, large image area, optional caption."""
        slide = self._add_slide()
        BackgroundManager.apply_solid(slide, FourthBrand.WHITE)

        # Title
        title_h = Inches(0.9)
        tf_title = self._add_textbox(
            slide,
            FourthBrand.MARGIN_LEFT,
            FourthBrand.MARGIN_TOP,
            FourthBrand.CONTENT_WIDTH,
            title_h,
        )
        TextFormatter.format_title(tf_title, title, level=2, color=FourthBrand.DEEP_BLUE)

        # Image area
        img_top = FourthBrand.MARGIN_TOP + title_h + Inches(0.2)
        caption_reserve = Inches(0.6) if caption else Inches(0)
        img_h = (
            FourthBrand.SLIDE_HEIGHT
            - img_top
            - FourthBrand.MARGIN_BOTTOM
            - caption_reserve
        )
        img_w = FourthBrand.CONTENT_WIDTH
        img_left = FourthBrand.MARGIN_LEFT

        inserted = None
        if image_path:
            inserted = ImageHandler.insert_from_file(
                slide, image_path, img_left, img_top, img_w, img_h,
            )
        elif image_b64:
            inserted = ImageHandler.insert_from_base64(
                slide, image_b64, img_left, img_top, img_w, img_h,
            )

        if inserted is None:
            ImageHandler.create_placeholder(
                slide, img_left, img_top, img_w, img_h, label="Image",
            )

        # Caption
        if caption:
            cap_top = img_top + img_h + Inches(0.08)
            tf_cap = self._add_textbox(
                slide, img_left, cap_top, img_w, Inches(0.5),
            )
            tf_cap.word_wrap = True
            p = tf_cap.paragraphs[0]
            p.text = caption
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.CAPTION_SIZE
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.DARK_GRAY
            TextFormatter._apply_line_spacing(p)

        return slide

    def add_quote_slide(
        self,
        quote: str,
        attribution: str,
        source_title: Optional[str] = None,
    ):
        """Quote slide: Teal accent bar on left, quote text, attribution below."""
        slide = self._add_slide()
        BackgroundManager.apply_solid(slide, FourthBrand.SOFT_WHITE)

        # Teal accent bar (vertical, left side)
        bar_left = Inches(1.5)
        bar_top = Inches(2.0)
        bar_h = Inches(3.0)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_left,
            bar_top,
            Inches(0.08),
            bar_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        bar.line.fill.background()

        # Quote text
        quote_left = bar_left + Inches(0.5)
        quote_w = Inches(9.5)
        quote_h = Inches(2.2)
        tf_q = self._add_textbox(slide, quote_left, bar_top + Inches(0.1), quote_w, quote_h)
        tf_q.word_wrap = True
        p = tf_q.paragraphs[0]
        # Use typographic open/close quotes
        p.text = f"\u201C{quote}\u201D"
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H3_SIZE
            run.font.bold = False
            run.font.italic = False  # Brand rule: never italic for emphasis
            run.font.color.rgb = FourthBrand.MIDNIGHT_NAVY
        TextFormatter._apply_line_spacing(p)

        # Attribution
        attr_top = bar_top + Inches(2.5)
        attr_h = Inches(0.8)
        tf_attr = self._add_textbox(slide, quote_left, attr_top, quote_w, attr_h)
        tf_attr.word_wrap = True
        p_attr = tf_attr.paragraphs[0]
        attr_text = f"-- {attribution}"
        if source_title:
            attr_text += f", {source_title}"
        p_attr.text = attr_text
        p_attr.alignment = PP_ALIGN.LEFT
        for run in p_attr.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.BODY_SMALL
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.DEEP_BLUE
        TextFormatter._apply_line_spacing(p_attr)

        return slide

    def add_data_slide(
        self,
        title: str,
        chart_placeholder_label: str = "Chart",
    ) -> Tuple[Any, Dict[str, Any]]:
        """Data slide with a placeholder area for a chart.

        Returns:
            (slide, chart_area) where *chart_area* is a dict with
            ``left``, ``top``, ``width``, ``height`` in EMU for chart placement.
        """
        slide = self._add_slide()
        BackgroundManager.apply_solid(slide, FourthBrand.WHITE)

        # Title
        title_h = Inches(0.9)
        tf_title = self._add_textbox(
            slide,
            FourthBrand.MARGIN_LEFT,
            FourthBrand.MARGIN_TOP,
            FourthBrand.CONTENT_WIDTH,
            title_h,
        )
        TextFormatter.format_title(tf_title, title, level=2, color=FourthBrand.DEEP_BLUE)

        # Chart placeholder area
        chart_top = FourthBrand.MARGIN_TOP + title_h + Inches(0.3)
        chart_h = FourthBrand.SLIDE_HEIGHT - chart_top - FourthBrand.MARGIN_BOTTOM
        chart_left = FourthBrand.MARGIN_LEFT
        chart_w = FourthBrand.CONTENT_WIDTH

        ImageHandler.create_placeholder(
            slide, chart_left, chart_top, chart_w, chart_h,
            label=chart_placeholder_label,
        )

        chart_area = {
            "left": chart_left,
            "top": chart_top,
            "width": chart_w,
            "height": chart_h,
        }
        return slide, chart_area

    def add_table_slide(
        self,
        title: str,
    ) -> Tuple[Any, Dict[str, Any]]:
        """Creates a slide with a title and returns coordinates for table placement.

        Returns:
            (slide, table_area) with ``left``, ``top``, ``width``, ``height``
            in EMU for placing a table via a companion TableBuilder.
        """
        slide = self._add_slide()
        BackgroundManager.apply_solid(slide, FourthBrand.WHITE)

        title_h = Inches(0.9)
        tf_title = self._add_textbox(
            slide,
            FourthBrand.MARGIN_LEFT,
            FourthBrand.MARGIN_TOP,
            FourthBrand.CONTENT_WIDTH,
            title_h,
        )
        TextFormatter.format_title(tf_title, title, level=2, color=FourthBrand.DEEP_BLUE)

        # Accent line
        line_top = FourthBrand.MARGIN_TOP + title_h + Inches(0.05)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            FourthBrand.MARGIN_LEFT,
            line_top,
            Inches(2.0),
            Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        accent.line.fill.background()

        table_top = line_top + Inches(0.3)
        table_h = FourthBrand.SLIDE_HEIGHT - table_top - FourthBrand.MARGIN_BOTTOM
        table_left = FourthBrand.MARGIN_LEFT
        table_w = FourthBrand.CONTENT_WIDTH

        table_area = {
            "left": table_left,
            "top": table_top,
            "width": table_w,
            "height": table_h,
        }
        return slide, table_area

    def add_closing_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        contact_info: Optional[str] = None,
        logo_path: Optional[str] = None,
    ):
        """Closing/thank-you slide: vignette background, white text, optional logo."""
        slide = self._add_slide()
        self._set_dark_bg(slide, 'vignette')

        # Title -- centred
        title_w = FourthBrand.CONTENT_WIDTH
        title_h = Inches(1.6)
        title_top = Inches(2.0) if subtitle or contact_info else Inches(2.8)
        tf = self._add_textbox(
            slide, self._centered_left(title_w), title_top, title_w, title_h,
        )
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H1_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        # Teal accent line below title
        accent_w = Inches(3.0)
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._centered_left(accent_w),
            title_top + title_h + Inches(0.1),
            accent_w, Inches(0.04),
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        accent_line.line.fill.background()

        current_top = title_top + title_h + Inches(0.25)

        # Subtitle
        if subtitle:
            sub_h = Inches(0.8)
            tf_sub = self._add_textbox(
                slide, self._centered_left(title_w), current_top, title_w, sub_h,
            )
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.alignment = PP_ALIGN.CENTER
            for run in p_sub.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.H3_SIZE
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE
            TextFormatter._apply_line_spacing(p_sub)
            current_top += sub_h + Inches(0.1)

        # Contact info
        if contact_info:
            ci_h = Inches(0.6)
            tf_ci = self._add_textbox(
                slide, self._centered_left(title_w), current_top, title_w, ci_h,
            )
            tf_ci.word_wrap = True
            p_ci = tf_ci.paragraphs[0]
            p_ci.text = contact_info
            p_ci.alignment = PP_ALIGN.CENTER
            for run in p_ci.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.BODY_SMALL
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE
            TextFormatter._apply_line_spacing(p_ci)

        # "Powered by iQ" floating badge at bottom-right
        self._add_floating_badge(
            slide, Inches(10.5), Inches(6.8),
            "Powered by iQ", bg_color='00B69F', text_color='FFFFFF', font_size_pt=9,
        )

        # Logo
        if logo_path:
            ImageHandler.insert_logo(slide, logo_path, position="bottom-center")

        return slide

    # -- v2 slide methods (design-first) --------------------------------- #

    def add_kpi_slide(
        self,
        metric_value: str,
        metric_label: str,
        context: Optional[str] = None,
        bg: str = "dark",
        category_label: Optional[str] = None,
    ):
        """Impact Moment KPI: 144pt hero number with optional progress ring.

        Radial gradient background, no glass card -- raw number dominates.
        If metric_value is a percentage, renders a donut ring visualization.

        Args:
            metric_value: The number to display (e.g., "15%", "$2.1M").
            metric_label: Short descriptor (e.g., "Uplift in Sales Per Labor Hour").
            context: Optional context line (e.g., "vs. 8% industry average").
            bg: "dark" (default) or "gradient" (vignette).
            category_label: Optional uppercase label (e.g., "IMPACT").
        """
        slide = self._add_slide()

        # -- Rich radial background (highlight behind the number) ----------
        OoxmlEffects.set_slide_gradient_bg(slide, stops=[
            (0, FourthBrand.HEX_DEEP_BLUE, 100),
            (40, FourthBrand.HEX_NAVY, 100),
            (100, '000D1A', 100),
        ], gradient_type='radial', center=(50, 30))

        # -- Category label (eyebrow) -------------------------------------
        if category_label:
            self.add_category_label(slide, category_label, color=FourthBrand.TEAL_GREEN)

        # -- Hero number at 144pt -----------------------------------------
        val_w = FourthBrand.CONTENT_WIDTH
        val_h = Inches(2.4)
        val_top = Inches(0.8) if category_label else Inches(0.6)

        # Use add_textbox directly so we can capture the shape for glow
        hero_box = slide.shapes.add_textbox(
            self._centered_left(val_w), val_top, val_w, val_h,
        )
        tf_val = hero_box.text_frame
        tf_val.word_wrap = True
        p = tf_val.paragraphs[0]
        p.text = str(metric_value)
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(144)
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE

        # Subtle glow on the hero number shape
        OoxmlEffects.add_glow(
            hero_box, radius_pt=20,
            hex_color=FourthBrand.HEX_TEAL, alpha_pct=15,
        )

        # -- Metric label at 28pt with letter spacing ---------------------
        label_w = Inches(8)
        label_h = Inches(0.8)
        label_top = val_top + val_h + Inches(0.1)
        label_box = slide.shapes.add_textbox(
            self._centered_left(label_w), label_top, label_w, label_h,
        )
        tf_label = label_box.text_frame
        tf_label.word_wrap = True
        p_l = tf_label.paragraphs[0]
        p_l.text = metric_label
        p_l.alignment = PP_ALIGN.CENTER
        for run in p_l.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(28)
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = FourthBrand.COOL_GREY
        # Subtle letter tracking
        if p_l.runs:
            OoxmlEffects.set_letter_spacing(p_l.runs[0], spacing_hundredths=150)

        # -- Progress ring or comparison bar ------------------------------
        # Try to parse metric_value as a percentage
        value_pct = None
        try:
            cleaned = str(metric_value).strip().rstrip('%')
            value_pct = float(cleaned)
            # Only treat as percentage if original had % or value is 0-100
            if '%' not in str(metric_value):
                value_pct = None
        except (ValueError, TypeError):
            value_pct = None

        if value_pct is not None:
            # Clamp to 0-100 for ring rendering
            ring_pct = max(0.0, min(100.0, value_pct))

            # Ring center position
            cx = FourthBrand.SLIDE_WIDTH // 2
            cy = Inches(4.8)
            outer_r = Inches(1.3)
            inner_r = Inches(0.95)

            # Background track ring (full circle, very faint white)
            OoxmlEffects.create_ring_segment(
                slide, cx, cy, outer_r, inner_r,
                start_deg=0, sweep_deg=360,
                fill_color='FFFFFF', alpha=8,
            )

            # Fill ring (partial, starting from top = 270 degrees)
            if ring_pct > 0:
                OoxmlEffects.create_ring_segment(
                    slide, cx, cy, outer_r, inner_r,
                    start_deg=270, sweep_deg=ring_pct * 3.6,
                    fill_color=FourthBrand.HEX_TEAL, alpha=90,
                )
        else:
            # Non-percentage: add a horizontal comparison bar
            bar_w = Inches(6)
            bar_left = self._centered_left(bar_w)
            bar_top = Inches(4.6)
            bar_h = Inches(0.12)

            # Track (faint white)
            track = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_left, bar_top, bar_w, bar_h,
            )
            OoxmlEffects.set_semi_transparent_fill(track, 'FFFFFF', 8)
            OoxmlEffects.set_rounded_corners(track, radius_pct=50)
            OoxmlEffects.remove_line(track)

            # Fill at 60% as a visual anchor
            fill_w = int(bar_w * 0.6)
            fill_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_left, bar_top, fill_w, bar_h,
            )
            fill_shape.fill.solid()
            fill_shape.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
            OoxmlEffects.set_rounded_corners(fill_shape, radius_pct=50)
            OoxmlEffects.remove_line(fill_shape)

        # -- Context line -------------------------------------------------
        if context:
            ctx_h = Inches(0.5)
            ctx_top = Inches(6.5)
            tf_ctx = self._add_textbox(
                slide, self._centered_left(val_w), ctx_top, val_w, ctx_h,
            )
            tf_ctx.word_wrap = True
            p_c = tf_ctx.paragraphs[0]
            p_c.text = context
            p_c.alignment = PP_ALIGN.CENTER
            for run in p_c.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(16)
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.COOL_GREY

        # -- "Powered by iQ" badge ----------------------------------------
        self._add_floating_badge(
            slide, Inches(10.5), Inches(6.8),
            "Powered by iQ",
            bg_color='00B69F', text_color='FFFFFF', font_size_pt=9,
        )

        return slide

    def add_kpi_grid(
        self,
        metrics: List[Dict[str, str]],
        bg: str = "dark",
    ):
        """KPI card grid. 2x2 or 1x3 layout depending on count.

        Args:
            metrics: List of dicts with 'value', 'label', optional 'change'.
                     e.g., [{'value': '23%', 'label': 'Cost Reduction', 'change': '+5%'}]
            bg: "dark" (default) or "gradient".
        """
        slide = self._add_slide()

        if bg == "gradient":
            self._set_dark_bg(slide, 'vignette')
        else:
            self._set_dark_bg(slide, 'auto')

        count = min(len(metrics), 4)
        if count <= 3:
            # 1x3 row layout
            cols, rows_count = count, 1
        else:
            # 2x2 grid
            cols, rows_count = 2, 2

        card_w = Inches(4.5)
        card_h = Inches(2.2)
        gutter_x = Inches(0.8)
        gutter_y = Inches(0.6)

        # Calculate starting positions
        total_w = cols * card_w + (cols - 1) * gutter_x
        total_h = rows_count * card_h + (rows_count - 1) * gutter_y
        start_left = (FourthBrand.SLIDE_WIDTH - total_w) // 2
        start_top = (FourthBrand.SLIDE_HEIGHT - total_h) // 2

        for idx, metric in enumerate(metrics[:count]):
            row = idx // cols
            col = idx % cols
            left = start_left + col * (card_w + gutter_x)
            top = start_top + row * (card_h + gutter_y)

            # Determine status color from change value
            change = metric.get("change", "")
            status_color = None
            if change:
                if str(change).startswith("-"):
                    status_color = FourthBrand.STATUS_ACTION
                elif str(change).startswith("+") or change:
                    status_color = FourthBrand.STATUS_GOOD

            # Glass-morphism card with status-colored border
            self._add_glass_card(
                slide, left, top, card_w, card_h,
                status_color=status_color or FourthBrand.STATUS_GOOD,
                glow=True,
            )

            # Value
            val_h = Inches(1.2)
            tf_val = self._add_textbox(
                slide, left + Inches(0.3), top + Inches(0.3),
                card_w - Inches(0.6), val_h,
            )
            tf_val.word_wrap = True
            p = tf_val.paragraphs[0]
            p.text = str(metric.get("value", ""))
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.H1_SIZE  # 48pt for grid cards
                run.font.bold = True
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE

            # Label
            lbl_h = Inches(0.5)
            lbl_top = top + Inches(0.3) + val_h
            tf_lbl = self._add_textbox(
                slide, left + Inches(0.3), lbl_top,
                card_w - Inches(0.6), lbl_h,
            )
            tf_lbl.word_wrap = True
            p_l = tf_lbl.paragraphs[0]
            p_l.text = str(metric.get("label", ""))
            p_l.alignment = PP_ALIGN.CENTER
            for run in p_l.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.BODY_SMALL
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.COOL_GREY

            # Change indicator (optional)
            change = metric.get("change")
            if change:
                chg_h = Inches(0.4)
                chg_top = lbl_top + lbl_h
                tf_chg = self._add_textbox(
                    slide, left + Inches(0.3), chg_top,
                    card_w - Inches(0.6), chg_h,
                )
                tf_chg.word_wrap = True
                p_c = tf_chg.paragraphs[0]
                p_c.text = str(change)
                p_c.alignment = PP_ALIGN.CENTER
                is_negative = str(change).startswith("-")
                for run in p_c.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = FourthBrand.CAPTION_SIZE
                    run.font.bold = True
                    run.font.italic = False
                    run.font.color.rgb = (
                        FourthBrand.HOT_RED if is_negative else FourthBrand.TEAL_GREEN
                    )

        return slide

    def add_pull_quote(
        self,
        quote: str,
        attribution: str,
        role: Optional[str] = None,
    ):
        """Pull-quote slide on dark background with teal accent bar.

        Large quote text on Midnight Navy background. Distinct from the lighter
        add_quote_slide() method.

        Args:
            quote: The quote text (do not include quotation marks).
            attribution: Person's name.
            role: Optional role/company (e.g., "Operations Director, Coastal Dining").
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Teal accent bar (vertical, left side) with glow
        bar_left = Inches(1.2)
        bar_top = Inches(1.8)
        bar_h = Inches(3.6)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_left,
            bar_top,
            Inches(0.06),
            bar_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        bar.line.fill.background()
        OoxmlEffects.add_glow(bar, radius_pt=10, hex_color=FourthBrand.HEX_TEAL,
                              alpha_pct=30)

        # Quote text
        quote_left = bar_left + Inches(0.5)
        quote_w = Inches(9.5)
        quote_h = Inches(2.8)
        tf_q = self._add_textbox(
            slide, quote_left, bar_top + Inches(0.2), quote_w, quote_h,
        )
        tf_q.word_wrap = True
        p = tf_q.paragraphs[0]
        p.text = f"\u201C{quote}\u201D"
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(28)
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        # Attribution
        attr_top = bar_top + Inches(3.2)
        attr_h = Inches(0.8)
        tf_attr = self._add_textbox(slide, quote_left, attr_top, quote_w, attr_h)
        tf_attr.word_wrap = True
        p_attr = tf_attr.paragraphs[0]
        p_attr.text = f"-- {attribution}"
        p_attr.alignment = PP_ALIGN.LEFT
        for run in p_attr.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.BODY_SMALL
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.TEAL_GREEN
        TextFormatter._apply_line_spacing(p_attr)

        # Role/company (optional)
        if role:
            role_top = attr_top + Inches(0.4)
            tf_role = self._add_textbox(
                slide, quote_left, role_top, quote_w, Inches(0.4),
            )
            tf_role.word_wrap = True
            p_r = tf_role.paragraphs[0]
            p_r.text = role
            p_r.alignment = PP_ALIGN.LEFT
            for run in p_r.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.BODY_SMALL
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.COOL_GREY

        return slide

    def add_problem_slide(self, statement: str):
        """Single powerful statement on dark background. 36-48pt white text centered.

        Used for problem statements, key claims, provocative questions.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        stmt_w = Inches(10.0)
        stmt_h = Inches(3.0)
        stmt_top = Inches(2.2)
        tf = self._add_textbox(
            slide, self._centered_left(stmt_w), stmt_top, stmt_w, stmt_h,
        )
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = statement
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H2_SIZE  # 36pt
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        return slide

    def add_comparison_slide(
        self,
        title: str,
        before: Union[str, List[str]],
        after: Union[str, List[str]],
    ):
        """Before/after two-column contrast slide.

        Left column is muted ("before" state), right column is teal-highlighted
        ("after" state). Visual contrast is the point.

        Args:
            title: Slide title (spans full width).
            before: String or bullet list for the "before" column.
            after: String or bullet list for the "after" column.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Title
        title_h = Inches(1.0)
        tf_title = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT, FourthBrand.MARGIN_TOP,
            FourthBrand.CONTENT_WIDTH, title_h,
        )
        TextFormatter.format_title(tf_title, title, level=2, color=FourthBrand.WHITE)

        # Accent line with glow
        line_top = FourthBrand.MARGIN_TOP + title_h + Inches(0.05)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, FourthBrand.MARGIN_LEFT, line_top,
            Inches(2.0), Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        accent.line.fill.background()
        OoxmlEffects.add_glow(accent, radius_pt=8, hex_color=FourthBrand.HEX_TEAL,
                              alpha_pct=25)

        # Column layout
        col_top = line_top + Inches(0.5)
        col_w = Inches(5.2)
        col_h = FourthBrand.SLIDE_HEIGHT - col_top - FourthBrand.MARGIN_BOTTOM
        gutter = Inches(1.33)
        left_left = FourthBrand.MARGIN_LEFT
        right_left = left_left + col_w + gutter

        # Glass card backgrounds for before/after columns
        card_padding = Inches(0.3)
        self._add_glass_card(
            slide, left_left - card_padding, col_top - card_padding,
            col_w + card_padding * 2, col_h + card_padding,
            status_color=FourthBrand.HEX_DARK_GRAY,
        )
        self._add_glass_card(
            slide, right_left - card_padding, col_top - card_padding,
            col_w + card_padding * 2, col_h + card_padding,
            status_color=FourthBrand.STATUS_GOOD,
            glow=True,
        )

        # "Before" header
        hdr_h = Inches(0.5)
        tf_bh = self._add_textbox(slide, left_left, col_top, col_w, hdr_h)
        tf_bh.word_wrap = True
        p = tf_bh.paragraphs[0]
        p.text = "BEFORE"
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.CATEGORY_LABEL_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.COOL_GREY

        # "After" header
        tf_ah = self._add_textbox(slide, right_left, col_top, col_w, hdr_h)
        tf_ah.word_wrap = True
        p_a = tf_ah.paragraphs[0]
        p_a.text = "AFTER"
        p_a.alignment = PP_ALIGN.LEFT
        for run in p_a.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.CATEGORY_LABEL_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.TEAL_GREEN

        # Column content
        content_top = col_top + hdr_h + Inches(0.15)
        content_h = col_h - hdr_h - Inches(0.15)

        for content, c_left, color in [
            (before, left_left, FourthBrand.COOL_GREY),
            (after, right_left, FourthBrand.WHITE),
        ]:
            tf = self._add_textbox(slide, c_left, content_top, col_w, content_h)
            # Handle dict input: extract 'points' list or 'title' string
            if isinstance(content, dict):
                items = content.get('points', content.get('items', []))
                if items:
                    TextFormatter.format_bullets(tf, items, color=color)
                else:
                    TextFormatter.format_body(tf, content.get('title', str(content)), color=color)
            elif isinstance(content, list):
                TextFormatter.format_bullets(tf, content, color=color)
            else:
                TextFormatter.format_body(tf, content, color=color)

        return slide

    @staticmethod
    def add_category_label(
        slide,
        label_text: str,
        color: Optional[RGBColor] = None,
    ):
        """Add small uppercase teal category label above slide title.

        12pt Poppins SemiBold, Teal color. E.g., 'WORKFORCE INTELLIGENCE'.
        Positioned at the standard margin top position.
        """
        color = color or FourthBrand.TEAL_GREEN

        label_h = Inches(0.3)
        txBox = slide.shapes.add_textbox(
            FourthBrand.MARGIN_LEFT,
            FourthBrand.MARGIN_TOP,
            Inches(6.0),
            label_h,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label_text.upper()
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.CATEGORY_LABEL_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = color

    def add_blank_slide(self):
        """Add a completely blank slide for custom content."""
        slide = self._add_slide()
        BackgroundManager.apply_solid(slide, FourthBrand.WHITE)
        return slide

    # -- v3 EBR-specific slide methods ----------------------------------- #

    def add_dark_content_slide(
        self,
        label: str,
        headline: str,
        body: Optional[Union[str, List[str]]] = None,
        sidebar_cards: Optional[List[Dict[str, str]]] = None,
    ):
        """Dark navy content slide with eyebrow label, headline, body text,
        and optional glass sidebar cards.

        Args:
            label: Uppercase eyebrow label (e.g., 'WORKFORCE INTELLIGENCE').
            headline: Slide headline in white.
            body: Body text string or bullet list.
            sidebar_cards: Optional list of dicts with 'title' and 'body'
                          rendered as glass cards on the right side.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Eyebrow label
        self.add_category_label(slide, label, color=FourthBrand.TEAL_GREEN)

        # Headline (reduced height to avoid accent line overlap)
        hl_top = FourthBrand.MARGIN_TOP + Inches(0.4)
        hl_h = Inches(1.0)
        body_w = FourthBrand.CONTENT_WIDTH
        if sidebar_cards:
            body_w = Inches(7.0)

        tf_hl = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT, hl_top, body_w, hl_h,
        )
        tf_hl.word_wrap = True
        p = tf_hl.paragraphs[0]
        p.text = headline
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
            OoxmlEffects.set_letter_spacing(run, 150)
        TextFormatter._apply_line_spacing(p)

        # Teal accent line (wider for more visual weight)
        line_top = hl_top + hl_h + Inches(0.05)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, FourthBrand.MARGIN_LEFT, line_top,
            Inches(3.0), Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        accent.line.fill.background()
        OoxmlEffects.add_glow(accent, radius_pt=6, hex_color=FourthBrand.HEX_TEAL,
                              alpha_pct=20)

        # Body content
        body_top = line_top + Inches(0.35)
        body_h = FourthBrand.SLIDE_HEIGHT - body_top - FourthBrand.MARGIN_BOTTOM

        if body:
            tf_body = self._add_textbox(
                slide, FourthBrand.MARGIN_LEFT, body_top, body_w, body_h,
            )
            if isinstance(body, list):
                TextFormatter.format_bullets(tf_body, body, color=FourthBrand.COOL_GREY)
            else:
                TextFormatter.format_body(tf_body, body, color=FourthBrand.COOL_GREY)

        # Sidebar glass cards — staggered sizes with decorative diagonal line
        if sidebar_cards:
            # Decorative diagonal line behind cards (rendered first so cards overlay)
            OoxmlEffects.add_decorative_line(
                slide,
                start_left=FourthBrand.MARGIN_LEFT + Inches(7.2),
                start_top=body_top,
                end_left=FourthBrand.MARGIN_LEFT + Inches(11.5),
                end_top=FourthBrand.SLIDE_HEIGHT - Inches(0.5),
                color='FFFFFF', alpha=8, width_pt=0.75,
            )

            card_left = FourthBrand.MARGIN_LEFT + Inches(7.5)
            card_w = Inches(4.0)
            n_cards = min(len(sidebar_cards), 3)
            card_gap = Inches(0.15)
            available_h = FourthBrand.SLIDE_HEIGHT - body_top - FourthBrand.MARGIN_BOTTOM
            usable_h = available_h - (n_cards - 1) * card_gap

            # First card gets 40% of space, remaining split the rest
            first_h_pct = 0.40
            remaining_pct = (1.0 - first_h_pct) / max(n_cards - 1, 1)

            for idx, card_data in enumerate(sidebar_cards[:3]):
                if idx == 0:
                    card_h_each = usable_h * first_h_pct
                else:
                    card_h_each = usable_h * remaining_pct

                if idx == 0:
                    card_top_pos = body_top
                else:
                    card_top_pos = body_top + usable_h * first_h_pct + card_gap
                    if idx > 1:
                        card_top_pos += (idx - 1) * (usable_h * remaining_pct + card_gap)
                self._add_glass_card(
                    slide, card_left, card_top_pos, card_w, card_h_each,
                    status_color=FourthBrand.STATUS_GOOD,
                )
                # Card title
                ct_h = Inches(0.4)
                tf_ct = self._add_textbox(
                    slide, card_left + Inches(0.2), card_top_pos + Inches(0.15),
                    card_w - Inches(0.4), ct_h,
                )
                tf_ct.word_wrap = True
                p_ct = tf_ct.paragraphs[0]
                p_ct.text = card_data.get('title', '')
                p_ct.alignment = PP_ALIGN.LEFT
                for run in p_ct.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = FourthBrand.BODY_SMALL
                    run.font.bold = True
                    run.font.italic = False
                    run.font.color.rgb = FourthBrand.WHITE
                # Card body
                tf_cb = self._add_textbox(
                    slide, card_left + Inches(0.2), card_top_pos + Inches(0.55),
                    card_w - Inches(0.4), card_h_each - Inches(0.7),
                )
                tf_cb.word_wrap = True
                p_cb = tf_cb.paragraphs[0]
                p_cb.text = card_data.get('body', '')
                p_cb.alignment = PP_ALIGN.LEFT
                for run in p_cb.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = FourthBrand.CAPTION_SIZE
                    run.font.bold = False
                    run.font.italic = False
                    run.font.color.rgb = FourthBrand.COOL_GREY

        return slide

    def add_stat_row_slide(
        self,
        label: str,
        headline: str,
        stats: List[Dict[str, str]],
        sparkline_data: Optional[List[List[float]]] = None,
    ):
        """Metric Dashboard: status-differentiated stat cards in a row.

        Each card uses color-coded fills and accent bars based on status:
        good=teal, watch=amber, action=red.

        Args:
            label: Uppercase eyebrow label.
            headline: Slide headline.
            stats: List of dicts with 'value', 'label', optional 'change', 'status'.
                   status: 'good' (teal), 'watch' (amber), 'action' (red).
            sparkline_data: Optional list of value lists (one per stat card)
                            for mini trend sparklines. Each inner list needs 2+
                            numeric values.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Eyebrow + headline
        self.add_category_label(slide, label, color=FourthBrand.TEAL_GREEN)
        hl_top = FourthBrand.MARGIN_TOP + Inches(0.4)
        hl_h = Inches(1.0)
        tf_hl = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT, hl_top,
            FourthBrand.CONTENT_WIDTH, hl_h,
        )
        tf_hl.word_wrap = True
        p = tf_hl.paragraphs[0]
        p.text = headline
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H2_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        # -- Status-differentiated card rendering -------------------------
        _STATUS_COLORS = {
            'good':   ('00B69F', 'FFFFFF'),   # (accent, card_fill_base)
            'watch':  ('FFB700', 'FFB700'),
            'action': ('D9373B', 'D9373B'),
        }

        count = min(len(stats), 4)
        card_w = Inches(2.6)
        card_h = Inches(3.2)
        gutter = Inches(0.4)
        total_w = count * card_w + (count - 1) * gutter
        start_left = (FourthBrand.SLIDE_WIDTH - total_w) // 2
        card_top = Inches(3.0)

        for idx, stat in enumerate(stats[:count]):
            left = start_left + idx * (card_w + gutter)
            status = stat.get('status', 'good')
            accent_hex, fill_hex = _STATUS_COLORS.get(status, ('00B69F', 'FFFFFF'))

            # -- Main card body (rounded, status-tinted fill) -------------
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_w, card_h,
            )
            OoxmlEffects.set_semi_transparent_fill(card, fill_hex, 6)
            OoxmlEffects.set_rounded_corners(card, radius_pct=5)
            OoxmlEffects.remove_line(card)

            # -- Top accent bar (full width, 3pt / 0.04") -----------------
            accent_h = Inches(0.04)
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, card_top, card_w, accent_h,
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(
                int(accent_hex[0:2], 16),
                int(accent_hex[2:4], 16),
                int(accent_hex[4:6], 16),
            )
            accent.line.fill.background()

            # -- Value text at 52pt ---------------------------------------
            val_text = str(stat.get('value', ''))
            val_h = Inches(1.2)
            tf_val = self._add_textbox(
                slide, left + Inches(0.1), card_top + Inches(0.25),
                card_w - Inches(0.2), val_h,
            )
            tf_val.word_wrap = True
            p_v = tf_val.paragraphs[0]
            p_v.text = val_text
            p_v.alignment = PP_ALIGN.CENTER
            for run in p_v.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(52)
                run.font.bold = True
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE

            # -- Label at 12pt -------------------------------------------
            lbl_h = Inches(0.5)
            lbl_top = card_top + Inches(1.45)
            tf_lbl = self._add_textbox(
                slide, left + Inches(0.1), lbl_top,
                card_w - Inches(0.2), lbl_h,
            )
            tf_lbl.word_wrap = True
            p_l = tf_lbl.paragraphs[0]
            p_l.text = str(stat.get('label', ''))
            p_l.alignment = PP_ALIGN.CENTER
            for run in p_l.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.CAPTION_SIZE
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.COOL_GREY

            # -- Change indicator with directional arrow ------------------
            change = stat.get('change')
            if change:
                change_str = str(change)
                chg_h = Inches(0.4)
                chg_top = lbl_top + lbl_h + Inches(0.05)
                tf_chg = self._add_textbox(
                    slide, left + Inches(0.1), chg_top,
                    card_w - Inches(0.2), chg_h,
                )
                tf_chg.word_wrap = True
                p_c = tf_chg.paragraphs[0]
                p_c.alignment = PP_ALIGN.CENTER

                # Determine direction and color
                is_neg = change_str.lstrip().startswith('-')
                if is_neg:
                    arrow_prefix = '\u25bc '   # down triangle
                    chg_color = FourthBrand.HOT_RED
                else:
                    arrow_prefix = '\u25b2 '   # up triangle
                    chg_color = FourthBrand.TEAL_GREEN

                p_c.text = arrow_prefix + change_str
                for run in p_c.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = FourthBrand.CAPTION_SIZE
                    run.font.bold = True
                    run.font.italic = False
                    run.font.color.rgb = chg_color

            # -- Mini sparkline (if data provided) ------------------------
            if sparkline_data and idx < len(sparkline_data):
                spark_values = sparkline_data[idx]
                if spark_values and len(spark_values) >= 2:
                    self._add_mini_sparkline(
                        slide,
                        left + Inches(0.3),
                        card_top + card_h - Inches(0.5),
                        card_w - Inches(0.6),
                        Inches(0.35),
                        spark_values,
                        color=accent_hex,
                    )

        return slide

    def add_outcome_slide(
        self,
        label: str,
        stat_value: str,
        stat_label: str,
        headline: str,
        body: Optional[str] = None,
        why_it_matters: Optional[str] = None,
        chart_area: bool = True,
        chart_data: Optional[Dict] = None,
    ) -> Union[Any, Tuple[Any, Dict[str, Any]]]:
        """3-zone outcome slide: info panel (stat+headline), full-width chart, callout.

        Zone A (1.1"-2.9"): Left 40% stat value/label, Right 60% headline/body.
        Zone B (3.0"-5.85"): Full-width gradient area chart.
        Zone C (6.0"-6.7"): Callout strip with teal accent.

        Args:
            label: Uppercase eyebrow label.
            stat_value: Large KPI number (e.g., '+23%'). Max 6 chars.
            stat_label: Label below the KPI. Max 25 chars.
            headline: Section headline. Max 50 chars (2 lines at 36pt).
            body: Optional body text below headline. Max 120 chars.
            why_it_matters: Optional callout strip text at bottom. Max 100 chars.
            chart_area: If True, returns (slide, chart_position_dict)
                        for ChartBuilder to add a chart.
            chart_data: Optional dict with 'categories', 'series', and
                optional 'title'. When provided, an area chart is rendered
                inline and the method returns only the slide.

        Returns:
            slide or (slide, chart_area_dict) if chart_area=True.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Eyebrow label
        self.add_category_label(slide, label, color=FourthBrand.TEAL_GREEN)

        # ── Zone A: Info panel (1.1" - 2.9") ──────────────────────────
        # Left 40%: stat value + label | Right 60%: headline + body
        left_w = Inches(4.5)
        left_top = Inches(1.1)
        right_left = FourthBrand.MARGIN_LEFT + left_w + Inches(0.5)
        right_w = FourthBrand.CONTENT_WIDTH - left_w - Inches(0.5)

        # Stat value (96pt KPI_LARGE_SIZE, left column)
        tf_val = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT,
            left_top, left_w, Inches(1.1),
        )
        tf_val.word_wrap = True
        p_v = tf_val.paragraphs[0]
        p_v.text = stat_value
        p_v.alignment = PP_ALIGN.LEFT
        for run in p_v.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.KPI_LARGE_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE

        # Stat label (below value, with letter spacing)
        tf_sl = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT,
            left_top + Inches(1.2), left_w, Inches(0.6),
        )
        tf_sl.word_wrap = True
        p_sl = tf_sl.paragraphs[0]
        p_sl.text = stat_label
        p_sl.alignment = PP_ALIGN.LEFT
        for run in p_sl.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.BODY_SIZE
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = FourthBrand.COOL_GREY
            OoxmlEffects.set_letter_spacing(run, 100)

        # Headline (right column, room for 2-line 36pt)
        tf_rh = self._add_textbox(
            slide, right_left, left_top, right_w, Inches(1.3),
        )
        tf_rh.word_wrap = True
        p_rh = tf_rh.paragraphs[0]
        p_rh.text = headline
        p_rh.alignment = PP_ALIGN.LEFT
        for run in p_rh.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H2_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE

        # Body text (below headline, only if provided)
        if body:
            body_size = FourthBrand.BODY_SMALL if chart_area else FourthBrand.BODY_SIZE
            tf_rb = self._add_textbox(
                slide, right_left, left_top + Inches(1.4),
                right_w, Inches(0.5),
            )
            TextFormatter.format_body(tf_rb, body, color=FourthBrand.COOL_GREY,
                                      size=body_size)

        # ── Divider line between Zone A and Zone B ─────────────────────
        self.add_divider_line(slide, FourthBrand.MARGIN_LEFT, Inches(2.85),
                              FourthBrand.CONTENT_WIDTH, color_hex='00B69F', alpha_pct=30)

        # ── Zone B: Chart (full width, 3.0" - 5.85") ─────────────────
        chart_top = Inches(3.0) if body else Inches(2.5)
        callout_top = Inches(6.0)
        chart_h = callout_top - chart_top - Inches(0.15)
        chart_pos = {
            'left': FourthBrand.MARGIN_LEFT,
            'top': chart_top,
            'width': FourthBrand.CONTENT_WIDTH,
            'height': chart_h,
        }

        # ── Zone C: Callout strip (6.0" - 6.7") ─────────────────────
        if why_it_matters:
            self.add_callout_strip(
                slide, why_it_matters,
                top=Inches(6.0),
            )

        # Inline chart rendering (when chart_data provided)
        if chart_data and chart_pos:
            from fourth_pptx_data import ChartBuilder
            cb = ChartBuilder()
            # Normalize series: accept list of dicts or list of tuples
            raw_series = chart_data['series']
            series = [
                (s['name'], s['values']) if isinstance(s, dict) else s
                for s in raw_series
            ]
            cb.add_dark_area_chart(
                slide,
                chart_data['categories'],
                series,
                title=chart_data.get('title'),
                position=chart_pos,
            )
            # Return slide only (chart already placed)
            return slide

        if chart_area:
            return slide, chart_pos
        return slide

    def add_gap_slide(
        self,
        label: str,
        headline: str,
        gaps: List[Dict[str, str]],
    ):
        """Card-per-row gap analysis slide with progress bars and status badges.

        Each gap is rendered as a data card row with area name, current/target
        values, a visual progress bar, and floating gap + status badges.

        Args:
            label: Uppercase eyebrow label.
            headline: Slide headline.
            gaps: List of dicts with 'area', 'current', 'best_practice',
                  'gap', 'status'. status: 'good'/'watch'/'action'.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Eyebrow + headline
        self.add_category_label(slide, label, color=FourthBrand.TEAL_GREEN)
        hl_top = FourthBrand.MARGIN_TOP + Inches(0.4)
        tf_hl = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT, hl_top,
            FourthBrand.CONTENT_WIDTH, Inches(0.8),
        )
        tf_hl.word_wrap = True
        p = tf_hl.paragraphs[0]
        p.text = headline
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H2_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        # Status color mapping (hex strings for OOXML calls)
        _STATUS_HEX = {
            'good': '00B69F',
            'watch': 'FFB700',
            'action': 'D9373B',
        }

        # Card-per-row layout
        row_top = Inches(2.2)
        row_h = Inches(1.05)
        row_gap_spacing = Inches(0.15)

        for idx, gap_data in enumerate(gaps[:6]):
            y = row_top + idx * (row_h + row_gap_spacing)
            area = gap_data.get('area', '')
            current = gap_data.get('current', '')
            best = gap_data.get('best_practice', '')
            gap_val = gap_data.get('gap', '')
            status = gap_data.get('status', 'watch')

            status_hex = _STATUS_HEX.get(status, 'FFB700')

            # Data card background (rectangle, no rounded corners)
            card = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, FourthBrand.MARGIN_LEFT, y,
                FourthBrand.CONTENT_WIDTH, row_h,
            )
            OoxmlEffects.set_semi_transparent_fill(card, 'FFFFFF', 4)
            OoxmlEffects.remove_line(card)

            # Bottom border line (subtle, status-colored)
            self.add_divider_line(
                slide, FourthBrand.MARGIN_LEFT, y + row_h,
                FourthBrand.CONTENT_WIDTH, color_hex=status_hex, alpha_pct=25,
            )

            # Area name (bold, 16pt, white)
            tf_area = self._add_textbox(
                slide, FourthBrand.MARGIN_LEFT + Inches(0.3), y + Inches(0.12),
                Inches(3.0), Inches(0.35),
            )
            tf_area.word_wrap = True
            p_a = tf_area.paragraphs[0]
            p_a.text = area
            p_a.alignment = PP_ALIGN.LEFT
            for run in p_a.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = FourthBrand.WHITE

            # "Current: XX%  |  Target: YY%" small text below area name
            tf_cur = self._add_textbox(
                slide, FourthBrand.MARGIN_LEFT + Inches(0.3), y + Inches(0.48),
                Inches(3.0), Inches(0.3),
            )
            tf_cur.word_wrap = True
            p_c = tf_cur.paragraphs[0]
            p_c.text = f"Current: {current}  |  Target: {best}"
            p_c.alignment = PP_ALIGN.LEFT
            for run in p_c.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(11)
                run.font.color.rgb = FourthBrand.COOL_GREY

            # Parse percentage for progress bar width
            try:
                cur_pct = float(current.replace('%', ''))
                bar_pct = min(cur_pct, 100)
            except (ValueError, AttributeError):
                bar_pct = 50  # fallback

            # Progress bar (center area)
            bar_left = FourthBrand.MARGIN_LEFT + Inches(3.8)
            bar_top = y + Inches(0.35)
            bar_w = Inches(5.5)
            bar_h = Inches(0.18)
            self._add_progress_bar_shape(
                slide, bar_left, bar_top, bar_w, bar_pct,
                height=bar_h, fill_color=status_hex, track_color='FFFFFF',
            )

            # Gap value as floating badge (right side)
            self._add_floating_badge(
                slide, FourthBrand.MARGIN_LEFT + Inches(10.0), y + Inches(0.15),
                gap_val, bg_color=status_hex, text_color='FFFFFF',
                font_size_pt=11,
            )

            # Status label badge below gap badge
            status_label = status.upper()
            self._add_floating_badge(
                slide, FourthBrand.MARGIN_LEFT + Inches(10.0), y + Inches(0.55),
                status_label, bg_color=status_hex, text_color='FFFFFF',
                font_size_pt=9,
            )

        return slide

    def add_recommendation_slide(
        self,
        label: str,
        headline: str,
        actions: List[Dict[str, str]],
    ):
        """Numbered action rows with timeline badges on dark background.

        Args:
            label: Uppercase eyebrow label.
            headline: Slide headline.
            actions: List of dicts with 'title', 'description',
                     optional 'timeline' (badge text like 'Q1 2026').
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Eyebrow + headline
        self.add_category_label(slide, label, color=FourthBrand.TEAL_GREEN)
        hl_top = FourthBrand.MARGIN_TOP + Inches(0.4)
        tf_hl = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT, hl_top,
            FourthBrand.CONTENT_WIDTH, Inches(0.8),
        )
        tf_hl.word_wrap = True
        p = tf_hl.paragraphs[0]
        p.text = headline
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H2_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        # Action rows
        row_top = Inches(2.4)
        row_h = Inches(1.1)
        row_gap = Inches(0.2)

        for idx, action in enumerate(actions[:5]):
            y = row_top + idx * (row_h + row_gap)

            # Glass card background for the row
            self._add_glass_card(
                slide, FourthBrand.MARGIN_LEFT, y,
                FourthBrand.CONTENT_WIDTH, row_h,
            )

            # Number circle
            circle_size = Inches(0.5)
            circle_left = FourthBrand.MARGIN_LEFT + Inches(0.3)
            circle_top = y + (row_h - circle_size) // 2
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, circle_left, circle_top,
                circle_size, circle_size,
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
            circle.line.fill.background()
            # Number text in circle
            tf_num = circle.text_frame
            tf_num.word_wrap = False
            p_num = tf_num.paragraphs[0]
            p_num.text = str(idx + 1)
            p_num.alignment = PP_ALIGN.CENTER
            for run in p_num.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = FourthBrand.WHITE

            # Action title
            text_left = circle_left + circle_size + Inches(0.3)
            text_w = Inches(7.5)
            tf_title = self._add_textbox(
                slide, text_left, y + Inches(0.1), text_w, Inches(0.4),
            )
            tf_title.word_wrap = True
            p_t = tf_title.paragraphs[0]
            p_t.text = action.get('title', '')
            p_t.alignment = PP_ALIGN.LEFT
            for run in p_t.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.BODY_SIZE
                run.font.bold = True
                run.font.italic = False
                run.font.color.rgb = FourthBrand.WHITE

            # Action description
            desc = action.get('description', '')
            if desc:
                tf_desc = self._add_textbox(
                    slide, text_left, y + Inches(0.5), text_w, Inches(0.5),
                )
                tf_desc.word_wrap = True
                p_d = tf_desc.paragraphs[0]
                p_d.text = desc
                p_d.alignment = PP_ALIGN.LEFT
                for run in p_d.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = FourthBrand.BODY_SMALL
                    run.font.bold = False
                    run.font.italic = False
                    run.font.color.rgb = FourthBrand.COOL_GREY

            # Timeline badge (right side)
            timeline = action.get('timeline', '')
            if timeline:
                badge_w = Inches(1.3)
                badge_h = Inches(0.35)
                badge_left = (FourthBrand.MARGIN_LEFT + FourthBrand.CONTENT_WIDTH
                              - badge_w - Inches(0.3))
                badge_top = y + (row_h - badge_h) // 2
                badge = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    badge_left, badge_top, badge_w, badge_h,
                )
                OoxmlEffects.set_semi_transparent_fill(badge, FourthBrand.HEX_TEAL, 25)
                OoxmlEffects.set_rounded_corners(badge, radius_pct=50)
                OoxmlEffects.remove_line(badge)
                tf_badge = badge.text_frame
                tf_badge.word_wrap = False
                p_b = tf_badge.paragraphs[0]
                p_b.text = timeline
                p_b.alignment = PP_ALIGN.CENTER
                for run in p_b.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = Pt(10)
                    run.font.bold = True
                    run.font.color.rgb = FourthBrand.WHITE

        return slide

    def add_roadmap_slide(
        self,
        label: str = 'ROADMAP RELEVANCE',
        headline: str = "What's Coming That Matters to You",
        pills: Optional[List[Dict[str, str]]] = None,
        items: Optional[List[Dict[str, str]]] = None,
    ):
        """Dark slide with pill badges and left-accent roadmap cards.

        Args:
            label: Uppercase eyebrow label.
            headline: Slide headline.
            pills: Optional list of dicts with 'text' and optional
                'style' ('solid' or 'outline'). Rendered as a row of
                pill badges below the headline.
            items: Optional list of dicts with 'title', 'description',
                'relevant' (why it matters), and 'timeline' (badge text).
                Up to 4 items rendered as left-accent cards.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'auto')

        # Eyebrow label + headline
        self.add_category_label(slide, label, color=FourthBrand.TEAL_GREEN)
        hl_top = FourthBrand.MARGIN_TOP + Inches(0.4)
        tf_hl = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT, hl_top,
            FourthBrand.CONTENT_WIDTH, Inches(1.0),
        )
        tf_hl.word_wrap = True
        p = tf_hl.paragraphs[0]
        p.text = headline
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H2_SIZE
            run.font.bold = True
            run.font.color.rgb = FourthBrand.WHITE

        # Row of pill badges (if provided)
        pill_top = hl_top + Inches(1.0)
        if pills:
            pill_left = FourthBrand.MARGIN_LEFT
            for pill in pills:
                style = pill.get('style', 'solid')
                is_outline = (style == 'outline')
                badge = self.add_pill_badge(
                    slide, pill['text'],
                    pill_left, pill_top,
                    outline_only=is_outline,
                )
                # Advance left by badge width + gap
                pill_left += badge.width + Inches(0.15)

        # Stack of left-accent cards (up to 4 items)
        card_start_top = pill_top + Inches(0.55)
        card_height = Inches(1.2)
        card_gap = Inches(0.2)

        if items:
            for i, item in enumerate(items[:4]):
                card_top = card_start_top + i * (card_height + card_gap)
                card, bar = self._add_left_accent_card(
                    slide, FourthBrand.MARGIN_LEFT, card_top,
                    FourthBrand.CONTENT_WIDTH, card_height,
                )

                # Title (bold white)
                tf_title = self._add_textbox(
                    slide,
                    FourthBrand.MARGIN_LEFT + Inches(0.25),
                    card_top + Inches(0.15),
                    FourthBrand.CONTENT_WIDTH - Inches(2.0),
                    Inches(0.3),
                )
                p_t = tf_title.paragraphs[0]
                p_t.text = item.get('title', '')
                p_t.alignment = PP_ALIGN.LEFT
                for run in p_t.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = Pt(16)
                    run.font.bold = True
                    run.font.color.rgb = FourthBrand.WHITE

                # Description (cool grey)
                desc = item.get('description', '')
                if desc:
                    tf_desc = self._add_textbox(
                        slide,
                        FourthBrand.MARGIN_LEFT + Inches(0.25),
                        card_top + Inches(0.45),
                        FourthBrand.CONTENT_WIDTH - Inches(2.0),
                        Inches(0.35),
                    )
                    tf_desc.word_wrap = True
                    p_desc = tf_desc.paragraphs[0]
                    p_desc.text = desc
                    p_desc.alignment = PP_ALIGN.LEFT
                    for run in p_desc.runs:
                        run.font.name = FourthBrand.FONT_FAMILY
                        run.font.size = FourthBrand.BODY_SMALL
                        run.font.color.rgb = FourthBrand.COOL_GREY

                # "Relevant because:" line (teal)
                relevant = item.get('relevant', '')
                if relevant:
                    tf_rel = self._add_textbox(
                        slide,
                        FourthBrand.MARGIN_LEFT + Inches(0.25),
                        card_top + Inches(0.82),
                        FourthBrand.CONTENT_WIDTH - Inches(2.0),
                        Inches(0.25),
                    )
                    tf_rel.word_wrap = True
                    p_rel = tf_rel.paragraphs[0]
                    # "Relevant because:" in bold teal, rest in regular grey
                    run_label = p_rel.add_run()
                    run_label.text = 'Relevant because: '
                    run_label.font.name = FourthBrand.FONT_FAMILY
                    run_label.font.size = FourthBrand.CAPTION_SIZE
                    run_label.font.bold = True
                    run_label.font.color.rgb = FourthBrand.TEAL_GREEN
                    run_val = p_rel.add_run()
                    run_val.text = relevant
                    run_val.font.name = FourthBrand.FONT_FAMILY
                    run_val.font.size = FourthBrand.CAPTION_SIZE
                    run_val.font.color.rgb = FourthBrand.COOL_GREY

                # Timeline pill badge (right-aligned, outline style)
                timeline = item.get('timeline', '')
                if timeline:
                    self.add_pill_badge(
                        slide, timeline,
                        FourthBrand.MARGIN_LEFT + FourthBrand.CONTENT_WIDTH - Inches(1.3),
                        card_top + Inches(0.15),
                        width=Inches(1.1), height=Inches(0.28),
                        outline_only=True,
                        outline_color_hex=FourthBrand.HEX_TEAL,
                    )

        return slide

    def add_mutual_commitments_slide(
        self,
        label: str = 'NEXT STEPS',
        headline: str = 'Mutual Commitments',
        fourth_items: Optional[List[str]] = None,
        client_items: Optional[List[str]] = None,
        client_name: str = 'Client',
        next_review_date: Optional[str] = None,
        contact_name: Optional[str] = None,
        contact_email: Optional[str] = None,
    ):
        """Split-panel partnership commitments slide with accent strips.

        Two visually distinct columns (deep blue / sky blue tint) with
        checkmark circles and accent strip items. Compact next-review bar
        at the bottom guarantees the section stays on-slide.

        Args:
            label: Eyebrow label.
            headline: Slide headline.
            fourth_items: List of Fourth's commitments.
            client_items: List of client's commitments.
            client_name: Name for the client column header.
            next_review_date: Optional date string for next review section.
            contact_name: Optional contact name shown below the date.
            contact_email: Optional contact email shown below the date.
        """
        slide = self._add_slide()
        self._set_dark_bg(slide, 'vignette')

        # Eyebrow + headline (centered)
        self.add_category_label(slide, label, color=FourthBrand.TEAL_GREEN)
        hl_top = FourthBrand.MARGIN_TOP + Inches(0.4)
        tf_hl = self._add_textbox(
            slide, FourthBrand.MARGIN_LEFT, hl_top,
            FourthBrand.CONTENT_WIDTH, Inches(0.8),
        )
        tf_hl.word_wrap = True
        p = tf_hl.paragraphs[0]
        p.text = headline
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.H2_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        TextFormatter._apply_line_spacing(p)

        # Split-panel layout dimensions
        col_w = Inches(5.5)
        col_h = Inches(3.6)
        gutter = Inches(0.73)
        col_top = Inches(2.0)
        left_left = FourthBrand.MARGIN_LEFT
        right_left = left_left + col_w + gutter

        # ---- Left column (Fourth's commitments) ---- #

        # Panel background: deep blue tint
        left_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left_left, col_top, col_w, col_h,
        )
        OoxmlEffects.set_semi_transparent_fill(left_panel, '0C4A7D', 15)
        OoxmlEffects.remove_line(left_panel)

        # Teal accent bar on left edge
        left_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left_left, col_top,
            Inches(0.05), col_h,
        )
        left_accent.fill.solid()
        left_accent.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
        left_accent.line.fill.background()

        # Header: "FOURTH COMMITS TO"
        tf_fh = self._add_textbox(
            slide, left_left + Inches(0.25), col_top + Inches(0.15),
            col_w - Inches(0.5), Inches(0.35),
        )
        tf_fh.word_wrap = True
        p_fh = tf_fh.paragraphs[0]
        p_fh.text = 'FOURTH COMMITS TO'
        p_fh.alignment = PP_ALIGN.LEFT
        for run in p_fh.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.CATEGORY_LABEL_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.TEAL_GREEN
        if p_fh.runs:
            OoxmlEffects.set_letter_spacing(p_fh.runs[0], spacing_hundredths=200)

        # Fourth items with check circles + accent strips
        if fourth_items:
            item_top = col_top + Inches(0.65)
            check_size = Inches(0.22)
            for i, item in enumerate(fourth_items[:5]):
                item_y = item_top + Inches(i * 0.55)

                # Teal check circle
                check = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left_left + Inches(0.15), item_y + Inches(0.1),
                    check_size, check_size,
                )
                check.fill.solid()
                check.fill.fore_color.rgb = FourthBrand.TEAL_GREEN
                check.line.fill.background()
                tf_chk = check.text_frame
                tf_chk.word_wrap = False
                tf_chk.margin_left = 0
                tf_chk.margin_right = 0
                tf_chk.margin_top = 0
                tf_chk.margin_bottom = 0
                p_chk = tf_chk.paragraphs[0]
                p_chk.alignment = PP_ALIGN.CENTER
                run_chk = p_chk.add_run()
                run_chk.text = "\u2713"
                run_chk.font.name = FourthBrand.FONT_FAMILY
                run_chk.font.size = Pt(10)
                run_chk.font.bold = True
                run_chk.font.color.rgb = FourthBrand.WHITE

                # Accent strip (shifted right for check circle)
                self._add_accent_strip(
                    slide,
                    left_left + Inches(0.50), item_y,
                    col_w - Inches(0.65), item,
                    accent_color='00B69F', font_size_pt=13,
                )

        # ---- Right column (Client's commitments) ---- #

        # Panel background: sky blue tint
        right_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, right_left, col_top, col_w, col_h,
        )
        OoxmlEffects.set_semi_transparent_fill(right_panel, '6FB4E3', 8)
        OoxmlEffects.remove_line(right_panel)

        # Sky blue accent bar on left edge
        right_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, right_left, col_top,
            Inches(0.05), col_h,
        )
        right_accent.fill.solid()
        right_accent.fill.fore_color.rgb = FourthBrand.SKY_BLUE
        right_accent.line.fill.background()

        # Header: "{CLIENT} COMMITS TO"
        tf_ch = self._add_textbox(
            slide, right_left + Inches(0.25), col_top + Inches(0.15),
            col_w - Inches(0.5), Inches(0.35),
        )
        tf_ch.word_wrap = True
        p_ch = tf_ch.paragraphs[0]
        p_ch.text = f'{client_name.upper()} COMMITS TO'
        p_ch.alignment = PP_ALIGN.LEFT
        for run in p_ch.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.CATEGORY_LABEL_SIZE
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.SKY_BLUE
        if p_ch.runs:
            OoxmlEffects.set_letter_spacing(p_ch.runs[0], spacing_hundredths=200)

        # Client items with check circles + accent strips
        if client_items:
            item_top = col_top + Inches(0.65)
            check_size = Inches(0.22)
            for i, item in enumerate(client_items[:5]):
                item_y = item_top + Inches(i * 0.55)

                # Sky blue check circle
                check = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    right_left + Inches(0.15), item_y + Inches(0.1),
                    check_size, check_size,
                )
                check.fill.solid()
                check.fill.fore_color.rgb = FourthBrand.SKY_BLUE
                check.line.fill.background()
                tf_chk = check.text_frame
                tf_chk.word_wrap = False
                tf_chk.margin_left = 0
                tf_chk.margin_right = 0
                tf_chk.margin_top = 0
                tf_chk.margin_bottom = 0
                p_chk = tf_chk.paragraphs[0]
                p_chk.alignment = PP_ALIGN.CENTER
                run_chk = p_chk.add_run()
                run_chk.text = "\u2713"
                run_chk.font.name = FourthBrand.FONT_FAMILY
                run_chk.font.size = Pt(10)
                run_chk.font.bold = True
                run_chk.font.color.rgb = FourthBrand.WHITE

                # Accent strip (shifted right for check circle)
                self._add_accent_strip(
                    slide,
                    right_left + Inches(0.50), item_y,
                    col_w - Inches(0.65), item,
                    accent_color='6FB4E3', font_size_pt=13,
                )

        # ---- NEXT REVIEW section (compact bar at bottom) ---- #
        if next_review_date:
            nr_top = Inches(6.4)
            nr_h = Inches(0.85)

            # Background bar
            nr_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, FourthBrand.MARGIN_LEFT, nr_top,
                FourthBrand.CONTENT_WIDTH, nr_h,
            )
            OoxmlEffects.set_semi_transparent_fill(nr_bg, 'FFFFFF', 6)
            OoxmlEffects.remove_line(nr_bg)

            # "NEXT REVIEW" label (left side)
            tf_nrl = self._add_textbox(
                slide, FourthBrand.MARGIN_LEFT + Inches(0.2), nr_top + Inches(0.08),
                Inches(2.0), Inches(0.25),
            )
            tf_nrl.word_wrap = False
            p_nrl = tf_nrl.paragraphs[0]
            p_nrl.text = 'NEXT REVIEW'
            p_nrl.alignment = PP_ALIGN.LEFT
            for run in p_nrl.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = FourthBrand.CATEGORY_LABEL_SIZE
                run.font.bold = True
                run.font.color.rgb = FourthBrand.TEAL_GREEN

            # Date (28pt, bold, white)
            tf_date = self._add_textbox(
                slide, FourthBrand.MARGIN_LEFT + Inches(0.2), nr_top + Inches(0.32),
                Inches(6.0), Inches(0.5),
            )
            tf_date.word_wrap = False
            p_d = tf_date.paragraphs[0]
            p_d.text = next_review_date
            p_d.alignment = PP_ALIGN.LEFT
            for run in p_d.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = FourthBrand.WHITE

            # Contact info (right-aligned)
            if contact_name or contact_email:
                contact_text = ''
                if contact_name:
                    contact_text = contact_name
                if contact_email:
                    contact_text += f' \u00b7 {contact_email}' if contact_text else contact_email
                tf_con = self._add_textbox(
                    slide, FourthBrand.MARGIN_LEFT + Inches(7.0), nr_top + Inches(0.35),
                    Inches(4.5), Inches(0.4),
                )
                tf_con.word_wrap = False
                p_con = tf_con.paragraphs[0]
                p_con.text = contact_text
                p_con.alignment = PP_ALIGN.RIGHT
                for run in p_con.runs:
                    run.font.name = FourthBrand.FONT_FAMILY
                    run.font.size = Pt(13)
                    run.font.color.rgb = FourthBrand.TEAL_GREEN

        return slide

    # -- v3 utility helpers ----------------------------------------------- #

    def add_callout_strip(
        self,
        slide,
        text: str,
        top: Optional[Any] = None,
        accent_color: Optional[str] = None,
    ):
        """Add a teal-accented callout strip (insight box) to a slide.

        A horizontal bar with a left teal accent and glass background,
        typically placed below charts.

        Args:
            slide: Target slide.
            text: Callout text.
            top: Vertical position (EMU). Default: near bottom.
            accent_color: Hex color for left accent bar. Default: teal.
        """
        accent_hex = accent_color or FourthBrand.HEX_TEAL
        strip_top = top or Inches(6.2)
        strip_left = FourthBrand.MARGIN_LEFT
        strip_w = FourthBrand.CONTENT_WIDTH
        strip_h = Inches(0.7)

        # Glass background (no status_color -- left bar provides accent)
        self._add_glass_card(
            slide, strip_left, strip_top, strip_w, strip_h,
        )

        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            strip_left, strip_top, Inches(0.06), strip_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(
            int(accent_hex[0:2], 16),
            int(accent_hex[2:4], 16),
            int(accent_hex[4:6], 16),
        )
        bar.line.fill.background()
        OoxmlEffects.add_glow(bar, radius_pt=6, hex_color=accent_hex,
                              alpha_pct=25)

        # Text
        tf = self._add_textbox(
            slide, strip_left + Inches(0.3), strip_top + Inches(0.1),
            strip_w - Inches(0.6), strip_h - Inches(0.2),
        )
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = FourthBrand.BODY_SMALL
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE

        return slide

    # -- v4 premium helpers ------------------------------------------------ #

    def add_pill_badge(
        self,
        slide,
        text: str,
        left, top,
        width=None,
        height=None,
        bg_color_hex: str = '00B69F',
        text_color=None,
        outline_only: bool = False,
        outline_color_hex: str = None,
    ):
        """Add a modern rounded pill/tag badge to a slide.

        Two visual modes:

        * **Solid** (default): filled background with white text.
        * **Outline-only** (``outline_only=True``): transparent fill with a
          coloured border and teal text.

        Args:
            slide: Target slide.
            text: Badge label.
            left: Horizontal position (EMU or ``Inches``).
            top: Vertical position (EMU or ``Inches``).
            width: Badge width. Default ``Inches(1.2)``.
            height: Badge height. Default ``Inches(0.28)``.
            bg_color_hex: Six-char hex colour for solid fill. Default teal.
            text_color: ``RGBColor`` override for label text.
            outline_only: When *True*, use transparent fill + outline.
            outline_color_hex: Hex colour for outline stroke (defaults to
                *bg_color_hex* when not supplied).
        """
        width = width or Inches(1.2)
        height = height or Inches(0.28)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        OoxmlEffects.set_rounded_corners(shape, radius_pct=50)

        if outline_only:
            OoxmlEffects.set_semi_transparent_fill(shape, '000000', alpha_pct=0)
            OoxmlEffects.set_line(
                shape,
                outline_color_hex or bg_color_hex,
                alpha_pct=80,
                width_pt=1.5,
            )
            resolved_color = text_color or FourthBrand.TEAL_GREEN
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(
                int(bg_color_hex[0:2], 16),
                int(bg_color_hex[2:4], 16),
                int(bg_color_hex[4:6], 16),
            )
            OoxmlEffects.remove_line(shape)
            resolved_color = text_color or FourthBrand.WHITE

        # Text
        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FourthBrand.FONT_FAMILY
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = resolved_color

        return shape

    def add_divider_line(
        self,
        slide,
        left, top, width,
        color_hex: str = 'FFFFFF',
        alpha_pct: int = 30,
        thickness: float = 0.015,
    ):
        """Add a thin horizontal separator line to a slide.

        Args:
            slide: Target slide.
            left: Horizontal position (EMU or ``Inches``).
            top: Vertical position (EMU or ``Inches``).
            width: Line width (EMU or ``Inches``).
            color_hex: Six-char hex fill colour. Default white.
            alpha_pct: Fill opacity percentage. Default 30.
            thickness: Line height in inches. Default 0.015.
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Inches(thickness),
        )
        OoxmlEffects.set_semi_transparent_fill(shape, color_hex, alpha_pct)
        OoxmlEffects.remove_line(shape)
        return shape

    def _add_left_accent_card(
        self,
        slide,
        left, top, width, height,
        accent_color_hex: str = '00B69F',
        accent_width: float = 0.06,
        fill_color_hex: str = 'FFFFFF',
        fill_alpha: int = 8,
    ):
        """Add a card with a thick left-only accent border.

        Returns a ``(card, bar)`` tuple so callers can populate the card's
        text frame and style the accent bar independently.

        Args:
            slide: Target slide.
            left: Card left edge (EMU or ``Inches``).
            top: Card top edge (EMU or ``Inches``).
            width: Total card width including accent bar.
            height: Card height.
            accent_color_hex: Hex colour for the left accent bar.
            accent_width: Accent bar width in inches. Default 0.06.
            fill_color_hex: Hex colour for the card body fill.
            fill_alpha: Card body fill opacity percentage. Default 8.
        """
        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(accent_width), height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(
            int(accent_color_hex[0:2], 16),
            int(accent_color_hex[2:4], 16),
            int(accent_color_hex[4:6], 16),
        )
        bar.line.fill.background()
        OoxmlEffects.add_glow(
            bar, radius_pt=4, hex_color=accent_color_hex, alpha_pct=20,
        )

        # Card body
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + Inches(accent_width), top,
            width - Inches(accent_width), height,
        )
        OoxmlEffects.set_semi_transparent_fill(card, fill_color_hex, fill_alpha)
        OoxmlEffects.set_rounded_corners(card, radius_pct=5)
        OoxmlEffects.remove_line(card)

        return (card, bar)

    def add_checkmark_item(
        self,
        slide,
        text: str,
        left, top, width,
        height=None,
        check_color=None,
    ):
        """Add a standalone teal checkmark + text line.

        Args:
            slide: Target slide.
            text: Item text displayed after the checkmark.
            left: Horizontal position (EMU or ``Inches``).
            top: Vertical position (EMU or ``Inches``).
            width: Text box width.
            height: Text box height. Default ``Inches(0.35)``.
            check_color: ``RGBColor`` for the checkmark glyph.
                Default ``FourthBrand.TEAL_GREEN``.
        """
        height = height or Inches(0.50)
        check_color = check_color or FourthBrand.TEAL_GREEN

        # Warn if text is long enough to wrap at 14pt
        if len(text) > 50:
            warnings.warn(
                f"Checkmark item >50 chars ({len(text)}): text will wrap to 2 lines."
            )

        tf = self._add_textbox(slide, left, top, width, height)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        # Checkmark run
        run_check = p.add_run()
        run_check.text = "\u2713  "
        run_check.font.name = FourthBrand.FONT_FAMILY
        run_check.font.size = Pt(16)
        run_check.font.bold = True
        run_check.font.color.rgb = check_color

        # Text run
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = FourthBrand.FONT_FAMILY
        run_text.font.size = FourthBrand.BODY_SMALL
        run_text.font.color.rgb = FourthBrand.WHITE

        return tf

    # -- v5 new visual primitives ----------------------------------------- #

    def _add_floating_badge(self, slide, left, top, text,
                            bg_color='00B69F', text_color='FFFFFF',
                            font_size_pt=10):
        """Small pill-shaped badge (rounded rectangle) with solid fill.

        Width auto-calculated from text length.

        Args:
            slide: Target slide.
            left: Horizontal position (EMU or Inches).
            top: Vertical position (EMU or Inches).
            text: Badge label text.
            bg_color: Six-char hex for fill. Default teal.
            text_color: Six-char hex for text. Default white.
            font_size_pt: Text size in points. Default 10.

        Returns:
            The created shape.
        """
        # Auto-calculate width: ~0.12" per char + 0.3" padding
        char_width = Inches(0.12)
        padding = Inches(0.3)
        badge_w = int(len(text) * char_width + padding)
        badge_h = Inches(0.28)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, badge_w, badge_h,
        )
        # Pill shape: corner radius = height/2 => 50% for full rounding
        OoxmlEffects.set_rounded_corners(shape, radius_pct=50)

        # Solid fill
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            int(bg_color[0:2], 16),
            int(bg_color[2:4], 16),
            int(bg_color[4:6], 16),
        )
        OoxmlEffects.remove_line(shape)

        # Text
        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FourthBrand.FONT_FAMILY
        run.font.size = Pt(font_size_pt)
        run.font.bold = True
        run.font.color.rgb = RGBColor(
            int(text_color[0:2], 16),
            int(text_color[2:4], 16),
            int(text_color[4:6], 16),
        )

        return shape

    def _add_accent_strip(self, slide, left, top, width, text,
                          accent_color='00B69F', font_size_pt=14):
        """Full-width container with left accent bar.

        Args:
            slide: Target slide.
            left: Left position (EMU or Inches).
            top: Top position (EMU or Inches).
            width: Total width (EMU or Inches).
            text: Text content for the strip.
            accent_color: Hex color for left accent bar. Default teal.
            font_size_pt: Text size in points. Default 14.

        Returns:
            The container shape.
        """
        strip_h = Inches(0.45)
        accent_w = Inches(0.08)

        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, accent_w, strip_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(
            int(accent_color[0:2], 16),
            int(accent_color[2:4], 16),
            int(accent_color[4:6], 16),
        )
        bar.line.fill.background()

        # Main container (semi-transparent white, 6% opacity)
        container_left = left + accent_w
        container_w = width - accent_w
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, container_left, top, container_w, strip_h,
        )
        OoxmlEffects.set_semi_transparent_fill(container, 'FFFFFF', alpha_pct=6)
        OoxmlEffects.remove_line(container)

        # Text inside container
        text_left = container_left + Inches(0.15)
        text_w = container_w - Inches(0.3)
        tf = self._add_textbox(slide, text_left, top, text_w, strip_h)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(font_size_pt)
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE

        return container

    def _add_metric_slab(self, slide, left, top, width, height,
                         value, label, sublabel=None):
        """Premium stat container with gradient fill.

        Args:
            slide: Target slide.
            left: Left position (EMU or Inches).
            top: Top position (EMU or Inches).
            width: Slab width (EMU or Inches).
            height: Slab height (EMU or Inches).
            value: The metric value string (e.g., "23%").
            label: Label below the value.
            sublabel: Optional secondary label.

        Returns:
            The background slab shape.
        """
        # Background shape with diagonal gradient
        slab = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        OoxmlEffects.set_rounded_corners(slab, radius_pct=8)
        OoxmlEffects.set_diagonal_gradient(slab, [
            {'pos': 0, 'color': FourthBrand.HEX_TEAL, 'alpha': 100},
            {'pos': 100, 'color': FourthBrand.HEX_DEEP_BLUE, 'alpha': 100},
        ], angle_deg=135)
        OoxmlEffects.remove_line(slab)

        # Inner shadow for depth
        OoxmlEffects.add_inner_shadow(
            slab, blur_pt=5, dist_pt=2, dir_deg=225,
            hex_color='000000', alpha_pct=30,
        )

        # Value text (72pt, centered)
        val_h = Inches(1.5)
        val_top = top + Inches(0.3)
        tf_val = self._add_textbox(
            slide, left + Inches(0.2), val_top,
            width - Inches(0.4), val_h,
        )
        tf_val.word_wrap = True
        p_v = tf_val.paragraphs[0]
        p_v.text = str(value)
        p_v.alignment = PP_ALIGN.CENTER
        for run in p_v.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(72)
            run.font.bold = True
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE

        # Label text (16pt, white)
        lbl_h = Inches(0.5)
        lbl_top = val_top + val_h + Inches(0.05)
        tf_lbl = self._add_textbox(
            slide, left + Inches(0.2), lbl_top,
            width - Inches(0.4), lbl_h,
        )
        tf_lbl.word_wrap = True
        p_l = tf_lbl.paragraphs[0]
        p_l.text = label
        p_l.alignment = PP_ALIGN.CENTER
        for run in p_l.runs:
            run.font.name = FourthBrand.FONT_FAMILY
            run.font.size = Pt(16)
            run.font.bold = False
            run.font.italic = False
            run.font.color.rgb = FourthBrand.WHITE
        # Subtle letter spacing on label
        if p_l.runs:
            OoxmlEffects.set_letter_spacing(p_l.runs[0], spacing_hundredths=50)

        # Optional sublabel (12pt, cool grey)
        if sublabel:
            sub_h = Inches(0.4)
            sub_top = lbl_top + lbl_h + Inches(0.02)
            tf_sub = self._add_textbox(
                slide, left + Inches(0.2), sub_top,
                width - Inches(0.4), sub_h,
            )
            tf_sub.word_wrap = True
            p_s = tf_sub.paragraphs[0]
            p_s.text = sublabel
            p_s.alignment = PP_ALIGN.CENTER
            for run in p_s.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(12)
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = FourthBrand.COOL_GREY

        return slab

    def _add_progress_bar_shape(self, slide, left, top, width, current_pct,
                                height=None, track_color='FFFFFF',
                                fill_color='00B69F', label=None):
        """Wrapper around OoxmlEffects.create_progress_bar().

        Args:
            slide: Target slide.
            left: Left position (EMU or Inches).
            top: Top position (EMU or Inches).
            width: Total bar width (EMU or Inches).
            current_pct: Fill percentage 0-100.
            height: Bar height. Default Inches(0.12).
            track_color: Hex color for track. Default white.
            fill_color: Hex color for fill. Default teal.
            label: If truthy, add percentage label to the right.

        Returns:
            (track_shape, fill_shape) tuple, or
            (track_shape, fill_shape, label_shape) if label is provided.
        """
        bar_height = height or Inches(0.12)

        track, fill = OoxmlEffects.create_progress_bar(
            slide, left, top, width, bar_height,
            fill_pct=current_pct,
            track_color=track_color,
            fill_color=fill_color,
            track_alpha=20,
            corner_radius_pct=50,
        )

        if label:
            # Add percentage label to the right of the bar
            label_w = Inches(0.6)
            label_left = left + width + Inches(0.1)
            tf = self._add_textbox(
                slide, label_left, top, label_w, bar_height,
            )
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = f"{current_pct}%"
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = FourthBrand.FONT_FAMILY
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = FourthBrand.WHITE
            return (track, fill, tf)

        return (track, fill)

    def _add_mini_sparkline(self, slide, left, top, width, height,
                            values, color='00B69F'):
        """Tiny trend visualization built from a freeform shape.

        Normalizes values to fit within width x height, then draws
        connected line segments.

        Args:
            slide: Target slide.
            left: Left position (EMU or Inches).
            top: Top position (EMU or Inches).
            width: Sparkline width (EMU or Inches).
            height: Sparkline height (EMU or Inches).
            values: List of numeric values (at least 2).
            color: Hex color for the line. Default teal.

        Returns:
            The created freeform shape.
        """
        if len(values) < 2:
            return None

        # Normalize values to the coordinate space
        min_val = min(values)
        max_val = max(values)
        val_range = max_val - min_val or 1  # avoid division by zero

        n = len(values)
        step_x = width / (n - 1)

        # Build freeform
        # Start at the first point
        first_y = top + height - int(height * (values[0] - min_val) / val_range)
        first_x = left

        freeform_builder = slide.shapes.build_freeform(first_x, first_y)

        for i in range(1, n):
            px = left + int(step_x * i)
            py = top + height - int(height * (values[i] - min_val) / val_range)
            freeform_builder.add_line_segments([(px, py)])

        shape = freeform_builder.convert_to_shape()

        # Style: line only, no fill
        shape.fill.background()  # transparent fill

        # Set line color and width
        shape.line.color.rgb = RGBColor(
            int(color[0:2], 16),
            int(color[2:4], 16),
            int(color[4:6], 16),
        )
        shape.line.width = Pt(1.5)

        return shape

    # -- properties ------------------------------------------------------- #

    @property
    def presentation(self):
        """Access the underlying ``Presentation`` object for advanced usage."""
        return self._prs

    # -- save ------------------------------------------------------------- #

    def save(self, filename: str):
        """Save the presentation to *filename*."""
        self._prs.save(filename)


# ---------------------------------------------------------------------------
# Class 6: PresentationRebrander -- rebrand existing PPTX files
# ---------------------------------------------------------------------------

class PresentationRebrander:
    """Rebrands an existing PPTX to Fourth standards.

    Typical usage::

        rebrander = PresentationRebrander("old-deck.pptx")
        rebrander.rebrand("fourth-deck.pptx")
    """

    def __init__(self, source_path: str):
        """Load the presentation from *source_path*."""
        if not os.path.isfile(source_path):
            raise FileNotFoundError(f"Source presentation not found: {source_path}")
        self._prs = Presentation(source_path)
        self._source = source_path

    # -- content extraction ----------------------------------------------- #

    def extract_content(self) -> Dict[str, Any]:
        """Extract all text, images, and notes into a structured dict.

        Returns::

            {
                "slides": [
                    {
                        "index": 0,
                        "shapes": [
                            {"type": "text", "text": "...", "paragraphs": [...]},
                            {"type": "image", "name": "...", "width": ..., "height": ...},
                        ],
                        "notes": "...",
                    },
                    ...
                ]
            }
        """
        result: Dict[str, Any] = {"slides": []}
        for idx, slide in enumerate(self._prs.slides):
            slide_data: Dict[str, Any] = {
                "index": idx,
                "shapes": [],
                "notes": "",
            }

            # Notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                slide_data["notes"] = notes_tf.text if notes_tf else ""

            # Shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paras = []
                    for p in shape.text_frame.paragraphs:
                        paras.append(p.text)
                    slide_data["shapes"].append({
                        "type": "text",
                        "text": shape.text_frame.text,
                        "paragraphs": paras,
                    })
                elif shape.shape_type is not None and hasattr(shape, "image"):
                    try:
                        slide_data["shapes"].append({
                            "type": "image",
                            "name": shape.name,
                            "width": shape.width,
                            "height": shape.height,
                        })
                    except Exception:
                        pass

            result["slides"].append(slide_data)
        return result

    # -- colour matching -------------------------------------------------- #

    @staticmethod
    def _color_distance(c1: RGBColor, c2: RGBColor) -> float:
        """Euclidean distance in RGB space between two colours."""
        return math.sqrt(
            (int(c1[0]) - int(c2[0])) ** 2
            + (int(c1[1]) - int(c2[1])) ** 2
            + (int(c1[2]) - int(c2[2])) ** 2
        )

    @staticmethod
    def _nearest_brand_color(rgb: RGBColor) -> RGBColor:
        """Find the closest Fourth brand colour to *rgb* by Euclidean distance."""
        best = FourthBrand.ALL_COLORS[0]
        best_dist = PresentationRebrander._color_distance(rgb, best)
        for candidate in FourthBrand.ALL_COLORS[1:]:
            d = PresentationRebrander._color_distance(rgb, candidate)
            if d < best_dist:
                best = candidate
                best_dist = d
        return best

    # -- recolouring ------------------------------------------------------ #

    @staticmethod
    def _is_dark_text(rgb: RGBColor) -> bool:
        """Return True if the colour is dark enough to be considered body/heading text."""
        luminance = 0.299 * int(rgb[0]) + 0.587 * int(rgb[1]) + 0.114 * int(rgb[2])
        return luminance < 128

    @staticmethod
    def _is_light_text(rgb: RGBColor) -> bool:
        """Return True if the colour is light (likely text-on-dark-bg)."""
        luminance = 0.299 * int(rgb[0]) + 0.587 * int(rgb[1]) + 0.114 * int(rgb[2])
        return luminance >= 200

    def _recolor_text(self, shape):
        """Recolour text runs in *shape* to Fourth standards.

        - Large/bold text (likely headings) -> Deep Blue
        - Dark body text -> Midnight Navy (replaces pure black)
        - Light text on presumably dark bg -> keep white
        """
        if not shape.has_text_frame:
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # Update font family
                run.font.name = FourthBrand.FONT_FAMILY
                # Never use italic for emphasis (brand rule)
                run.font.italic = False

                # Determine if this is a heading or body
                font_size = run.font.size
                is_heading = (
                    run.font.bold
                    or (font_size is not None and font_size >= FourthBrand.H3_SIZE)
                )

                # Get current colour if set
                try:
                    current_rgb = run.font.color.rgb
                except (AttributeError, TypeError):
                    current_rgb = None

                if current_rgb is None:
                    # No explicit colour -- apply defaults
                    if is_heading:
                        run.font.color.rgb = FourthBrand.DEEP_BLUE
                    else:
                        run.font.color.rgb = FourthBrand.MIDNIGHT_NAVY
                elif self._is_light_text(current_rgb):
                    # Light text (white-ish) -- keep white
                    run.font.color.rgb = FourthBrand.WHITE
                elif self._is_dark_text(current_rgb):
                    if is_heading:
                        run.font.color.rgb = FourthBrand.DEEP_BLUE
                    else:
                        # Replace pure black with Midnight Navy
                        run.font.color.rgb = FourthBrand.MIDNIGHT_NAVY
                else:
                    # Mid-tone colour -- map to nearest brand colour
                    run.font.color.rgb = self._nearest_brand_color(current_rgb)

    def _recolor_fills(self, shape):
        """Recolour shape fills to the nearest Fourth brand colour."""
        try:
            fill = shape.fill
        except AttributeError:
            return

        # Only handle solid fills we can read
        try:
            if fill.type is not None and hasattr(fill, "fore_color"):
                current = fill.fore_color.rgb
                if current is not None:
                    fill.fore_color.rgb = self._nearest_brand_color(current)
        except (AttributeError, TypeError):
            pass

        # Line colour
        try:
            line = shape.line
            if line.color and line.color.rgb:
                line.color.rgb = self._nearest_brand_color(line.color.rgb)
        except (AttributeError, TypeError):
            pass

    # -- theme application ------------------------------------------------ #

    def apply_fourth_theme(self):
        """Apply Fourth branding to all slides.

        - Fonts -> Poppins
        - Dark text -> Midnight Navy / Deep Blue
        - Shape fills -> nearest brand colour
        - Remove italics (brand rule)
        """
        for slide in self._prs.slides:
            for shape in slide.shapes:
                self._recolor_text(shape)
                self._recolor_fills(shape)

                # Handle grouped shapes
                if shape.shape_type is not None and hasattr(shape, "shapes"):
                    try:
                        for child in shape.shapes:
                            self._recolor_text(child)
                            self._recolor_fills(child)
                    except Exception:
                        pass

    # -- full rebrand ----------------------------------------------------- #

    def rebrand(self, output_path: str):
        """Full rebrand: apply theme + save.

        This is the primary entry point.  It:
        1. Sets slide dimensions to 16:9
        2. Applies all Fourth colour and typography rules
        3. Saves to *output_path*
        """
        # Force 16:9 if not already
        self._prs.slide_width = FourthBrand.SLIDE_WIDTH
        self._prs.slide_height = FourthBrand.SLIDE_HEIGHT

        self.apply_fourth_theme()
        self._prs.save(output_path)

    @property
    def presentation(self):
        """Access the underlying ``Presentation`` for advanced modifications."""
        return self._prs


# ---------------------------------------------------------------------------
# Convenience: module-level quick-build function
# ---------------------------------------------------------------------------

def quick_deck(
    title: str,
    subtitle: str = "",
    slides: Optional[List[Dict[str, Any]]] = None,
    output: str = "presentation.pptx",
) -> str:
    """Build a simple Fourth-branded deck from minimal input.

    Each item in *slides* is a dict with at least ``"title"`` and one of
    ``"bullets"``, ``"body"``, ``"quote"``, or ``"image_path"``.

    Returns the output filename.
    """
    builder = PresentationBuilder()
    builder.add_title_slide(title, subtitle)

    for s in (slides or []):
        stype = s.get("type", "content")
        stitle = s.get("title", "")

        if stype == "section":
            builder.add_section_break(stitle, s.get("subtitle"))
        elif stype == "kpi":
            builder.add_kpi_slide(
                s.get("value", ""),
                s.get("label", stitle),
                context=s.get("context"),
                category_label=s.get("category_label"),
            )
        elif stype == "kpi_grid":
            builder.add_kpi_grid(s.get("metrics", []))
        elif stype == "pull_quote":
            builder.add_pull_quote(
                s.get("quote", ""),
                s.get("attribution", ""),
                s.get("role"),
            )
        elif stype == "problem":
            builder.add_problem_slide(s.get("statement", stitle))
        elif stype == "comparison":
            builder.add_comparison_slide(
                stitle,
                s.get("before", ""),
                s.get("after", ""),
            )
        elif stype == "quote":
            builder.add_quote_slide(
                s.get("quote", ""),
                s.get("attribution", ""),
                s.get("source_title"),
            )
        elif stype == "image":
            builder.add_image_slide(
                stitle,
                image_path=s.get("image_path"),
                image_b64=s.get("image_b64"),
                caption=s.get("caption"),
            )
        elif stype == "two_column":
            builder.add_two_column(
                stitle,
                s.get("left", ""),
                s.get("right", ""),
            )
        elif stype == "three_column":
            builder.add_three_column(stitle, s.get("columns", []))
        else:
            # Default: content slide
            builder.add_content_slide(
                stitle,
                bullets=s.get("bullets"),
                body_text=s.get("body"),
                layout=s.get("layout", "single"),
                category_label=s.get("category_label"),
            )

    builder.add_closing_slide("Thank You", subtitle="Powered by iQ")
    builder.save(output)
    return output
