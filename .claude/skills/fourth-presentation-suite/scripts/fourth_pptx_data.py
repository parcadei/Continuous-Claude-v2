"""
Fourth Presentation Suite -- Data Visualization
Charts, tables, and data formatting for Fourth-branded presentations.

Runs in Claude.ai sandbox (no pip, no network). Only python-pptx + stdlib.

Usage:
    from fourth_pptx_core import PresentationBuilder
    from fourth_pptx_data import ChartBuilder, TableBuilder, DataFormatter

    builder = PresentationBuilder()
    slide, area = builder.add_data_slide("Revenue by Quarter")

    chart_builder = ChartBuilder()
    chart_builder.add_bar_chart(slide,
        categories=["Q1", "Q2", "Q3", "Q4"],
        series=[("Revenue", [1.2, 1.5, 1.8, 2.1])],
        title="Quarterly Revenue ($M)", position=area)

    builder.save("revenue.pptx")
"""
from __future__ import annotations
from typing import Any, Callable, Dict, List, Optional, Tuple, Union

from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.oxml.ns import qn

from fourth_ooxml import OoxmlEffects

# -- Fourth Brand Colors (mirrors brand-essentials.md) ---------------------
DEEP_BLUE = RGBColor(0x0C, 0x4A, 0x7D)
TEAL_GREEN = RGBColor(0x00, 0xB6, 0x9F)
SKY_BLUE = RGBColor(0x6F, 0xB4, 0xE3)
MIDNIGHT_NAVY = RGBColor(0x00, 0x27, 0x47)
DARK_GRAY = RGBColor(0x37, 0x3E, 0x42)
COOL_GREY = RGBColor(0xCF, 0xD1, 0xD1)
SOFT_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HOT_RED = RGBColor(0xD9, 0x37, 0x3B)
PURPLE = RGBColor(0x92, 0x79, 0xB2)
BRAND_FONT = "Poppins"

# -- v3 Dark theme chart colors -----------------------------------------------
DARK_CHART_TEXT = RGBColor(0x7A, 0x95, 0xA8)    # muted text on dark bg
DARK_GRIDLINE = RGBColor(0x1A, 0x3A, 0x5C)      # near-invisible on navy

_DEFAULT_POS = {
    "left": Inches(1.0), "top": Inches(1.8),
    "width": Inches(11.0), "height": Inches(5.0),
}


# ======================================================================
# ChartBuilder
# ======================================================================
class ChartBuilder:
    """Fourth-branded charts on python-pptx slides.

    All add_* methods accept slide, data args, optional title, and optional
    position dict (left/top/width/height as Inches or Emu).  Brand styling
    is applied automatically.
    """

    CHART_COLORS: List[RGBColor] = [
        DEEP_BLUE, TEAL_GREEN, SKY_BLUE, PURPLE, HOT_RED,
    ]
    DARK_CHART_COLORS: List[RGBColor] = [
        TEAL_GREEN, SKY_BLUE,
        RGBColor(0x00, 0x8A, 0x7A),  # dark teal
        RGBColor(0x4E, 0xC9, 0xB0),  # light teal
        RGBColor(0x00, 0x6B, 0x5E),  # deep teal
    ]

    # -- Public: category-based charts ------------------------------------

    def add_bar_chart(self, slide, categories, series, title=None, position=None):
        """Clustered horizontal bar chart."""
        return self._add_category_chart(
            slide, XL_CHART_TYPE.BAR_CLUSTERED, categories, series, title, position, "bar")

    def add_column_chart(self, slide, categories, series, title=None, position=None):
        """Clustered vertical column chart."""
        return self._add_category_chart(
            slide, XL_CHART_TYPE.COLUMN_CLUSTERED, categories, series, title, position, "column")

    def add_line_chart(self, slide, categories, series, title=None, position=None):
        """Line chart with markers -- good for trends over time."""
        return self._add_category_chart(
            slide, XL_CHART_TYPE.LINE_MARKERS, categories, series, title, position, "line")

    def add_stacked_bar(self, slide, categories, series, title=None, position=None):
        """Stacked bar chart for composition / parts-of-whole."""
        return self._add_category_chart(
            slide, XL_CHART_TYPE.BAR_STACKED, categories, series, title, position, "stacked_bar")

    def add_area_chart(self, slide, categories, series, title=None, position=None):
        """Area chart for cumulative trends."""
        return self._add_category_chart(
            slide, XL_CHART_TYPE.AREA, categories, series, title, position, "area")

    # -- Public: single-series charts -------------------------------------

    def add_pie_chart(self, slide, categories, values, title=None, position=None):
        """Pie chart (single series). values is a flat list of numbers."""
        return self._add_single_series_chart(
            slide, XL_CHART_TYPE.PIE, categories, values, title, position, "pie")

    def add_donut_chart(self, slide, categories, values, title=None, position=None):
        """Donut chart -- pie with a hole."""
        return self._add_single_series_chart(
            slide, XL_CHART_TYPE.DOUGHNUT, categories, values, title, position, "donut")

    # -- Public: scatter ---------------------------------------------------

    def add_scatter_chart(self, slide, data_points, title=None, position=None):
        """Scatter/XY chart. data_points: [(name, [(x,y), ...]), ...]"""
        pos = position or _DEFAULT_POS
        chart_data = XyChartData()
        for name, points in data_points:
            s = chart_data.add_series(name)
            for x, y in points:
                s.add_data_point(x, y)
        shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            pos["left"], pos["top"], pos["width"], pos["height"], chart_data)
        if title:
            shape.chart.has_title = True
            shape.chart.chart_title.text_frame.text = title
        self._apply_fourth_style(shape.chart, "scatter")
        return shape

    # -- Public: combo chart -----------------------------------------------

    def add_combo_chart(self, slide, categories, bar_series, line_series,
                        title=None, position=None):
        """Column + line combination chart.

        bar_series rendered as columns; line_series converted to lines via XML
        patch (python-pptx has no native combo support).
        """
        pos = position or _DEFAULT_POS
        all_series = list(bar_series) + list(line_series)
        chart_data = self._build_category_data(categories, all_series)
        shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            pos["left"], pos["top"], pos["width"], pos["height"], chart_data)
        chart = shape.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        self._convert_trailing_series_to_line(chart, len(bar_series), len(line_series))
        self._apply_fourth_style(chart, "combo")
        return shape

    # -- Internal: chart creation helpers ----------------------------------

    def _add_category_chart(self, slide, chart_type_enum, categories, series,
                            title, position, style_key):
        """Shared builder for multi-series category charts."""
        pos = position or _DEFAULT_POS
        chart_data = self._build_category_data(categories, series)
        shape = slide.shapes.add_chart(
            chart_type_enum,
            pos["left"], pos["top"], pos["width"], pos["height"], chart_data)
        chart = shape.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        self._apply_fourth_style(chart, style_key)
        return shape

    def _add_single_series_chart(self, slide, chart_type_enum, categories, values,
                                  title, position, style_key):
        """Shared builder for pie/donut (single series, flat values list)."""
        pos = position or _DEFAULT_POS
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("Values", values)
        shape = slide.shapes.add_chart(
            chart_type_enum,
            pos["left"], pos["top"], pos["width"], pos["height"], chart_data)
        chart = shape.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        self._apply_fourth_style(chart, style_key)
        return shape

    @staticmethod
    def _build_category_data(categories, series):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series:
            chart_data.add_series(name, values)
        return chart_data

    def _get_default_position(self):
        return dict(_DEFAULT_POS)

    # -- Internal: Fourth brand styling ------------------------------------

    def _apply_fourth_style(self, chart, chart_type="bar"):
        """Apply Fourth brand identity: colors, fonts, axes, legend, borders."""
        is_point_chart = chart_type in ("pie", "donut")

        # Series / point coloring
        if is_point_chart:
            for idx in range(len(chart.series[0].points)):
                pt = chart.series[0].points[idx]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = self.CHART_COLORS[idx % len(self.CHART_COLORS)]
        else:
            for idx, s in enumerate(chart.series):
                color = self.CHART_COLORS[idx % len(self.CHART_COLORS)]
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = color
                if chart_type in ("line", "combo"):
                    s.format.line.color.rgb = color
                    s.format.line.width = Pt(2.5)
                    if hasattr(s, "marker"):
                        s.marker.style = 8  # circle
                        s.marker.size = 8
                        s.marker.format.fill.solid()
                        s.marker.format.fill.fore_color.rgb = color
                        s.marker.format.line.color.rgb = color

        # Title -- v2: 20pt Poppins SemiBold Deep Blue
        if chart.has_title:
            self._style_text_frame(chart.chart_title.text_frame,
                                   font_size=Pt(20), font_bold=True, font_color=DEEP_BLUE)

        # Legend -- v2: bottom, Poppins Regular 12pt
        if not (chart_type == "scatter" and len(chart.series) <= 1):
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            self._style_font(chart.legend.font, BRAND_FONT, Pt(12), MIDNIGHT_NAVY)

        # Axes (skip for pie/donut)
        if not is_point_chart:
            self._style_axes(chart)

        # Data labels for pie/donut
        if is_point_chart:
            plot = chart.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.name = BRAND_FONT
            dl.font.size = Pt(12)
            dl.font.color.rgb = MIDNIGHT_NAVY
            dl.number_format = "0%"
            dl.number_format_is_linked = False
            dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

    def _style_axes(self, chart):
        """Style category and value axes with Fourth brand fonts/colors.

        v2: 12pt axis labels, softer gridlines (Cool Grey at reduced weight),
        no chart border/frame.
        """
        try:
            ca = chart.category_axis
            self._style_font(ca.tick_labels.font, BRAND_FONT, Pt(12), MIDNIGHT_NAVY)
            ca.has_major_gridlines = False
            ca.has_minor_gridlines = False
            ca.format.line.color.rgb = COOL_GREY
            ca.format.line.width = Pt(0.25)
        except (ValueError, AttributeError):
            pass
        try:
            va = chart.value_axis
            self._style_font(va.tick_labels.font, BRAND_FONT, Pt(12), MIDNIGHT_NAVY)
            va.has_major_gridlines = True
            va.has_minor_gridlines = False
            va.major_gridlines.format.line.color.rgb = COOL_GREY
            va.major_gridlines.format.line.width = Pt(0.15)  # softer gridlines
            va.format.line.color.rgb = COOL_GREY
            va.format.line.width = Pt(0.25)
        except (ValueError, AttributeError):
            pass

    @staticmethod
    def _style_font(font, name, size, color, bold=False):
        font.name = name
        font.size = size
        font.color.rgb = color
        font.bold = bold

    @staticmethod
    def _style_text_frame(tf, font_name=BRAND_FONT, font_size=Pt(14),
                          font_bold=False, font_color=MIDNIGHT_NAVY, alignment=None):
        for para in tf.paragraphs:
            if alignment is not None:
                para.alignment = alignment
            for run in para.runs:
                run.font.name = font_name
                run.font.size = font_size
                run.font.bold = font_bold
                run.font.color.rgb = font_color

    # -- v3: Dark theme charts (for navy/radial gradient backgrounds) ------

    def add_dark_bar_chart(self, slide, categories, series, title=None,
                           position=None):
        """Bar chart styled for dark backgrounds.

        Uses muted axis text, near-invisible gridlines, and transparent
        plot area background.
        """
        shape = self.add_bar_chart(slide, categories, series, title, position)
        self._apply_dark_style(shape.chart)
        return shape

    def add_dark_column_chart(self, slide, categories, series, title=None,
                              position=None):
        """Column chart styled for dark backgrounds."""
        shape = self.add_column_chart(slide, categories, series, title, position)
        self._apply_dark_style(shape.chart)
        return shape

    def add_dark_line_chart(self, slide, categories, series, title=None,
                            position=None):
        """Line chart styled for dark backgrounds."""
        shape = self.add_line_chart(slide, categories, series, title, position)
        self._apply_dark_style(shape.chart)
        return shape

    def add_dark_donut(self, slide, categories, values, title=None,
                       position=None, cutout=78):
        """High-cutout donut for dark backgrounds.

        Matches the HTML health-score visualization with a large center
        area for the primary metric label.

        Args:
            cutout: Hole size percentage (0-90). Default 78 for KPI display.
        """
        shape = self.add_donut_chart(slide, categories, values, title, position)
        chart = shape.chart
        self._apply_dark_style(chart)

        # Set the donut hole size via XML
        plot_elem = chart._chartSpace.chart.plotArea
        donut_charts = plot_elem.findall(qn('c:doughnutChart'))
        for dc in donut_charts:
            hole_size = dc.find(qn('c:holeSize'))
            if hole_size is not None:
                hole_size.set('val', str(cutout))
            else:
                hs = etree.SubElement(dc, qn('c:holeSize'))
                hs.set('val', str(cutout))

        return shape

    def add_dark_area_chart(
        self, slide, categories, series, title=None, position=None,
        gradient_color: str = '00B69F',
        gradient_top_alpha: int = 60,
        smooth: bool = True,
    ):
        """Premium area chart with gradient fill for dark backgrounds.

        Creates an area chart with a signature gradient-under-curve effect:
        teal fill fading from semi-transparent at top to fully transparent
        at bottom, smooth curves, and circle markers with white outlines.

        Args:
            slide: Target slide object.
            categories: List of category labels (e.g. month names).
            series: List of (name, values) tuples.
            title: Optional chart title string.
            position: Optional dict with left/top/width/height keys.
            gradient_color: Hex color for gradient fill (default teal).
            gradient_top_alpha: Top-of-fill opacity percentage (0-100).
            smooth: Whether to apply smooth curve interpolation.
        """
        shape = self._add_category_chart(
            slide, XL_CHART_TYPE.AREA, categories, series,
            title, position, "area"
        )
        self._apply_dark_style(shape.chart)

        # -- Post-creation OOXML injection for premium gradient effect --
        chart_space = shape.chart._chartSpace
        plot_area = chart_space.chart.plotArea

        area_chart = plot_area.find(qn('c:areaChart'))
        if area_chart is None:
            return shape  # fallback: no injection if structure unexpected

        for ser in area_chart.findall(qn('c:ser')):
            # Gradient fill: color fading to transparent
            OoxmlEffects.set_series_gradient_fill(
                ser, gradient_color, gradient_top_alpha, 0
            )
            # Visible line on top of the area
            OoxmlEffects.set_series_line_style(ser, gradient_color, 2.5)
            # Smooth curves
            if smooth:
                OoxmlEffects.enable_smooth_lines(ser)
            # Circle markers with white outline
            OoxmlEffects.add_markers_to_series(
                ser, 'circle', 8, gradient_color, 'FFFFFF', 1.5
            )

        return shape

    def _apply_dark_style(self, chart):
        """Override chart styling for dark backgrounds.

        - Title text -> white
        - Axis labels -> muted color (DARK_CHART_TEXT) at Pt(10)
        - Gridlines -> DARK_GRIDLINE at thin weight
        - Legend -> muted color, Pt(10), positioned at top
        - Plot area + chart area fill -> transparent
        - Series recolored to teal palette (visible on navy)
        """
        # Title
        if chart.has_title:
            self._style_text_frame(chart.chart_title.text_frame,
                                   font_size=Pt(20), font_bold=True,
                                   font_color=WHITE)

        # Legend -- top position, smaller font
        if chart.has_legend:
            self._style_font(chart.legend.font, BRAND_FONT, Pt(10),
                             DARK_CHART_TEXT)
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False

        # Axes -- Pt(10) for cleaner look
        try:
            ca = chart.category_axis
            self._style_font(ca.tick_labels.font, BRAND_FONT, Pt(10),
                             DARK_CHART_TEXT)
            ca.format.line.color.rgb = DARK_GRIDLINE
            ca.format.line.width = Pt(0.15)
        except (ValueError, AttributeError):
            pass
        try:
            va = chart.value_axis
            self._style_font(va.tick_labels.font, BRAND_FONT, Pt(10),
                             DARK_CHART_TEXT)
            va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = DARK_GRIDLINE
            va.major_gridlines.format.line.width = Pt(0.15)
            va.format.line.color.rgb = DARK_GRIDLINE
            va.format.line.width = Pt(0.15)
        except (ValueError, AttributeError):
            pass

        # Make plot area transparent via XML
        plot_elem = chart._chartSpace.chart.plotArea
        spPr = plot_elem.find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(plot_elem, qn('c:spPr'))
        for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
            existing = spPr.find(qn(tag))
            if existing is not None:
                spPr.remove(existing)
        etree.SubElement(spPr, qn('a:noFill'))

        # Clear chart area frame to transparent
        try:
            chart.chart_area.format.fill.background()
            chart.chart_area.format.line.fill.background()
        except (AttributeError, ValueError):
            pass

        # Recolor series to dark-friendly teal palette
        try:
            plot = chart.plots[0]
            for idx, series in enumerate(plot.series):
                color = self.DARK_CHART_COLORS[idx % len(self.DARK_CHART_COLORS)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
                series.format.line.color.rgb = color
        except (IndexError, AttributeError, ValueError):
            pass

    # -- Internal: combo chart XML manipulation ----------------------------

    @staticmethod
    def _convert_trailing_series_to_line(chart, num_bar, num_line):
        """Patch chart XML to move trailing series into a <c:lineChart> overlay."""
        if num_line == 0:
            return
        plot_area = chart.element.find(qn("c:plotArea"))
        bar_elem = plot_area.find(qn("c:barChart")) if plot_area is not None else None
        if bar_elem is None:
            return
        all_ser = bar_elem.findall(qn("c:ser"))
        if len(all_ser) < num_bar + num_line:
            return

        # Pull line series out of bar chart
        line_sers = all_ser[num_bar:]
        for s in line_sers:
            bar_elem.remove(s)

        # Build <c:lineChart>
        from lxml import etree
        lc = etree.SubElement(plot_area, qn("c:lineChart"))
        etree.SubElement(lc, qn("c:grouping")).set("val", "standard")
        etree.SubElement(lc, qn("c:varyColors")).set("val", "0")
        for s in line_sers:
            lc.append(s)
        etree.SubElement(lc, qn("c:marker")).set("val", "1")

        # Secondary axis wiring
        bar_ax_ids = bar_elem.findall(qn("c:axId"))
        if len(bar_ax_ids) < 2:
            return
        pri_cat = bar_ax_ids[0].get("val")
        pri_val = bar_ax_ids[1].get("val")
        sec_cat = str(int(pri_cat) + 100)
        sec_val = str(int(pri_val) + 100)

        etree.SubElement(lc, qn("c:axId")).set("val", sec_cat)
        etree.SubElement(lc, qn("c:axId")).set("val", sec_val)

        # Secondary category axis (hidden, shared categories)
        sca = etree.SubElement(plot_area, qn("c:catAx"))
        etree.SubElement(sca, qn("c:axId")).set("val", sec_cat)
        sc = etree.SubElement(sca, qn("c:scaling"))
        etree.SubElement(sc, qn("c:orientation")).set("val", "minMax")
        etree.SubElement(sca, qn("c:delete")).set("val", "1")
        etree.SubElement(sca, qn("c:axPos")).set("val", "b")
        etree.SubElement(sca, qn("c:crossAx")).set("val", sec_val)

        # Secondary value axis (visible, right side)
        sva = etree.SubElement(plot_area, qn("c:valAx"))
        etree.SubElement(sva, qn("c:axId")).set("val", sec_val)
        sv = etree.SubElement(sva, qn("c:scaling"))
        etree.SubElement(sv, qn("c:orientation")).set("val", "minMax")
        etree.SubElement(sva, qn("c:delete")).set("val", "0")
        etree.SubElement(sva, qn("c:axPos")).set("val", "r")
        etree.SubElement(sva, qn("c:crossAx")).set("val", sec_cat)
        etree.SubElement(sva, qn("c:crosses")).set("val", "max")


# ======================================================================
# TableBuilder
# ======================================================================
class TableBuilder:
    """Fourth-branded tables: Deep Blue headers, alternating row shading,
    Poppins font, Midnight Navy text, auto-detected numeric alignment.
    """

    def add_table(self, slide, headers, rows, position=None, col_widths=None):
        """Create a styled table on the slide.

        Args:
            headers: Column header strings.
            rows: List of row-data lists (str/number values).
            position: Dict with left/top/width/height.
            col_widths: Optional list of column widths (Inches/Emu).
        Returns: the Table object.
        """
        pos = position or _DEFAULT_TABLE_POS
        n_rows, n_cols = len(rows) + 1, len(headers)
        shape = slide.shapes.add_table(
            n_rows, n_cols, pos["left"], pos["top"], pos["width"], pos["height"])
        table = shape.table

        # Column widths
        widths = col_widths or self._auto_col_widths(headers, rows, pos["width"])
        for i, w in enumerate(widths):
            if i < n_cols:
                table.columns[i].width = w

        # Populate cells
        for ci, h in enumerate(headers):
            table.cell(0, ci).text = str(h)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                if ci < n_cols:
                    table.cell(ri + 1, ci).text = str(val)

        self._style_header_row(table)
        self._style_data_rows(table)
        return table

    def add_kpi_table(self, slide, kpis, position=None):
        """KPI display table.

        kpis: list of dicts with 'label', 'value', optional 'change', 'unit'.
        Renders large bold values with color-coded change indicators.
        """
        has_change = any("change" in k for k in kpis)
        n_cols = 3 if has_change else 2
        headers = ["Metric", "Value"] + (["Change"] if has_change else [])
        pos = position or {
            "left": Inches(2.0), "top": Inches(2.0),
            "width": Inches(9.0), "height": Inches(0.6 * len(kpis) + 0.4),
        }
        n_rows = len(kpis) + 1
        shape = slide.shapes.add_table(
            n_rows, n_cols, pos["left"], pos["top"], pos["width"], pos["height"])
        table = shape.table

        # Column proportions
        tw = int(pos["width"]) if not isinstance(pos["width"], int) else pos["width"]
        if n_cols == 3:
            table.columns[0].width = int(tw * 0.45)
            table.columns[1].width = int(tw * 0.30)
            table.columns[2].width = int(tw * 0.25)
        else:
            table.columns[0].width = int(tw * 0.55)
            table.columns[1].width = int(tw * 0.45)

        # Header text
        for ci, h in enumerate(headers):
            table.cell(0, ci).text = h
        self._style_header_row(table)

        # KPI rows
        for ri, kpi in enumerate(kpis):
            dr = ri + 1
            bg = SOFT_WHITE if dr % 2 == 0 else WHITE
            unit = kpi.get("unit", "")
            val = kpi.get("value", "")

            table.cell(dr, 0).text = str(kpi.get("label", ""))
            table.cell(dr, 1).text = f"{unit}{val}" if unit else str(val)
            self._style_cell(table.cell(dr, 0), bg, MIDNIGHT_NAVY, Pt(14), False, PP_ALIGN.LEFT)
            self._style_cell(table.cell(dr, 1), bg, DEEP_BLUE, Pt(20), True, PP_ALIGN.CENTER)

            if has_change:
                ch = str(kpi.get("change", ""))
                table.cell(dr, 2).text = ch
                color = HOT_RED if ch.startswith("-") else (
                    TEAL_GREEN if ch and (ch[0] == "+" or ch[0].isdigit()) else MIDNIGHT_NAVY)
                self._style_cell(table.cell(dr, 2), bg, color, Pt(14), True, PP_ALIGN.CENTER)

        return table

    # -- Styling helpers ---------------------------------------------------

    def _style_header_row(self, table):
        """Deep Blue bg, White Poppins Semibold, center-aligned."""
        for ci in range(len(table.columns)):
            self._style_cell(table.cell(0, ci), DEEP_BLUE, WHITE, Pt(14), True, PP_ALIGN.CENTER)

    def _style_data_rows(self, table):
        """Alternating Soft White / White, Midnight Navy text, numbers right-aligned."""
        for ri in range(1, len(table.rows)):
            bg = SOFT_WHITE if ri % 2 == 0 else WHITE
            for ci in range(len(table.columns)):
                cell = table.cell(ri, ci)
                align = PP_ALIGN.RIGHT if _looks_numeric(cell.text.strip()) else PP_ALIGN.LEFT
                self._style_cell(cell, bg, MIDNIGHT_NAVY, Pt(14), False, align)

    @staticmethod
    def _style_cell(cell, bg_color, text_color, font_size=Pt(14),
                    bold=False, alignment=PP_ALIGN.LEFT):
        """Apply full styling to a single table cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        for para in cell.text_frame.paragraphs:
            para.alignment = alignment
            para.space_before = Pt(0)
            para.space_after = Pt(0)
            for run in para.runs:
                run.font.name = BRAND_FONT
                run.font.size = font_size
                run.font.bold = bold
                run.font.color.rgb = text_color
            # Fallback: paragraph-level font when no runs exist
            if not para.runs and para.text:
                para.font.name = BRAND_FONT
                para.font.size = font_size
                para.font.bold = bold
                para.font.color.rgb = text_color

    def add_matrix_table(self, slide, headers, rows, position=None):
        """Feature comparison matrix with check/cross marks.

        Args:
            headers: Column headers (e.g., ["Feature", "Fourth", "Competitor A"]).
            rows: List of row lists. Values can be:
                  True / "check" -> Teal checkmark
                  False / "cross" -> Cool Grey cross
                  "partial" -> Sky Blue partial indicator
                  Any other string -> displayed as-is.
            position: Optional dict with left/top/width/height.
        Returns: the Table object.
        """
        pos = position or _DEFAULT_TABLE_POS
        n_rows, n_cols = len(rows) + 1, len(headers)
        shape = slide.shapes.add_table(
            n_rows, n_cols, pos["left"], pos["top"], pos["width"], pos["height"])
        table = shape.table

        # Column widths: first column wider for feature names
        tw = int(pos["width"]) if not isinstance(pos["width"], int) else pos["width"]
        first_w = int(tw * 0.35)
        rest_w = int((tw - first_w) / max(n_cols - 1, 1))
        if n_cols > 0:
            table.columns[0].width = first_w
        for i in range(1, n_cols):
            table.columns[i].width = rest_w

        # Header row
        for ci, h in enumerate(headers):
            table.cell(0, ci).text = str(h)
        self._style_header_row(table)

        # Data rows with check/cross/partial indicators
        for ri, row in enumerate(rows):
            dr = ri + 1
            bg = SOFT_WHITE if dr % 2 == 0 else WHITE
            for ci, val in enumerate(row):
                if ci >= n_cols:
                    break
                cell = table.cell(dr, ci)

                if ci == 0:
                    # Feature name column
                    cell.text = str(val)
                    self._style_cell(cell, bg, MIDNIGHT_NAVY, Pt(14), False, PP_ALIGN.LEFT)
                else:
                    # Indicator columns
                    if val is True or str(val).lower() in ("check", "yes", "true"):
                        cell.text = "\u2713"  # checkmark
                        self._style_cell(cell, bg, TEAL_GREEN, Pt(18), True, PP_ALIGN.CENTER)
                    elif val is False or str(val).lower() in ("cross", "no", "false"):
                        cell.text = "\u2717"  # cross mark
                        self._style_cell(cell, bg, COOL_GREY, Pt(18), True, PP_ALIGN.CENTER)
                    elif str(val).lower() == "partial":
                        cell.text = "\u25D0"  # half circle
                        self._style_cell(cell, bg, SKY_BLUE, Pt(18), True, PP_ALIGN.CENTER)
                    else:
                        cell.text = str(val)
                        self._style_cell(cell, bg, MIDNIGHT_NAVY, Pt(14), False, PP_ALIGN.CENTER)

        return table

    def add_kpi_card(self, slide, value, label, position, change=None, color=None):
        """Styled KPI card shape: large bold number + label below.

        Args:
            slide: Target slide.
            value: The metric value string (e.g., "23%").
            label: The metric label (e.g., "Cost Reduction").
            position: Dict with left/top/width/height.
            change: Optional change indicator (e.g., "+5%"). Green if positive, red if negative.
            color: Override card background color (default Deep Blue).
        Returns: the card shape.
        """
        from pptx.enum.shapes import MSO_SHAPE

        card_color = color or DEEP_BLUE
        pos = position

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            pos["left"], pos["top"], pos["width"], pos["height"],
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.fill.background()

        # Value text
        val_h = int(pos["height"] * 0.5)
        txBox_val = slide.shapes.add_textbox(
            pos["left"] + Inches(0.2),
            pos["top"] + Inches(0.2),
            pos["width"] - Inches(0.4),
            val_h,
        )
        tf_val = txBox_val.text_frame
        tf_val.word_wrap = True
        p = tf_val.paragraphs[0]
        p.text = str(value)
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = BRAND_FONT
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = WHITE

        # Label text
        lbl_h = int(pos["height"] * 0.25)
        lbl_top = pos["top"] + Inches(0.2) + val_h
        txBox_lbl = slide.shapes.add_textbox(
            pos["left"] + Inches(0.2), lbl_top,
            pos["width"] - Inches(0.4), lbl_h,
        )
        tf_lbl = txBox_lbl.text_frame
        tf_lbl.word_wrap = True
        p_l = tf_lbl.paragraphs[0]
        p_l.text = str(label)
        p_l.alignment = PP_ALIGN.CENTER
        for run in p_l.runs:
            run.font.name = BRAND_FONT
            run.font.size = Pt(12)
            run.font.bold = False
            run.font.color.rgb = COOL_GREY

        # Change indicator
        if change:
            chg_top = lbl_top + lbl_h
            chg_h = int(pos["height"] * 0.2)
            txBox_chg = slide.shapes.add_textbox(
                pos["left"] + Inches(0.2), chg_top,
                pos["width"] - Inches(0.4), chg_h,
            )
            tf_chg = txBox_chg.text_frame
            tf_chg.word_wrap = True
            p_c = tf_chg.paragraphs[0]
            p_c.text = str(change)
            p_c.alignment = PP_ALIGN.CENTER
            is_negative = str(change).startswith("-")
            for run in p_c.runs:
                run.font.name = BRAND_FONT
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = HOT_RED if is_negative else TEAL_GREEN

        return card

    def add_dark_table(self, slide, headers, rows, position=None,
                       col_widths=None, modern=True):
        """Table styled for dark backgrounds using semi-transparent rows.

        Header: Deep Blue bg, white text. Data rows: semi-transparent
        glass-like fills with muted text colors.

        Args:
            modern: If True, apply borderless modern styling with teal
                    header accent line instead of grid borders.
        """
        table = self.add_table(slide, headers, rows, position, col_widths)
        self._style_dark_data_rows(table)
        if modern:
            self._apply_modern_style(table, dark=True)
        return table

    # OOXML spec: CT_TableCellProperties child order
    _TCPR_CHILD_ORDER = [
        'a:lnL', 'a:lnR', 'a:lnT', 'a:lnB', 'a:lnTlToBr', 'a:lnBlToTr',
        'a:cell3D',
        # fills (mutually exclusive)
        'a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill',
        'a:pattFill', 'a:grpFill',
        'a:headers',
    ]

    @staticmethod
    def _reorder_tcPr_children(tcPr):
        """Re-sort tcPr children to match OOXML required order.

        CRITICAL: Must be called after any XML injection into tcPr.
        PowerPoint will show a repair dialog if children are out of order.
        """
        children = list(tcPr)
        for child in children:
            tcPr.remove(child)

        def sort_key(elem):
            tag = etree.QName(elem.tag).localname
            ns = etree.QName(elem.tag).namespace or ''
            prefix = 'a' if 'drawingml' in ns else 'p'
            full_tag = f'{prefix}:{tag}'
            try:
                return TableBuilder._TCPR_CHILD_ORDER.index(full_tag)
            except ValueError:
                return len(TableBuilder._TCPR_CHILD_ORDER)

        children.sort(key=sort_key)
        for child in children:
            tcPr.append(child)

    @staticmethod
    def _remove_cell_borders(cell):
        """Remove all cell borders by setting them to noFill."""
        tc = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))
        for border_tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
            existing = tcPr.find(qn(border_tag))
            if existing is not None:
                tcPr.remove(existing)
            ln = etree.SubElement(tcPr, qn(border_tag))
            ln.set('w', '0')
            etree.SubElement(ln, qn('a:noFill'))

    @staticmethod
    def _add_header_accent_line(table, col_count):
        """Add a 2pt teal bottom border to header row cells only."""
        for ci in range(col_count):
            cell = table.cell(0, ci)
            tc = cell._tc
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            # Remove existing bottom border
            existing = tcPr.find(qn('a:lnB'))
            if existing is not None:
                tcPr.remove(existing)
            # Add 2pt teal bottom border
            lnB = etree.SubElement(tcPr, qn('a:lnB'))
            lnB.set('w', str(int(Pt(2))))
            solidFill = etree.SubElement(lnB, qn('a:solidFill'))
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgb.set('val', '00B69F')

    def _apply_modern_style(self, table, dark=False):
        """Remove grid borders, add teal header accent, apply subtle fills.

        Args:
            table: The python-pptx Table object.
            dark: If True, use semi-transparent white fills for dark bg.
                  If False, use light alternating fills.
        """
        n_rows = len(table.rows)
        n_cols = len(table.columns)

        # Remove all cell borders
        for ri in range(n_rows):
            for ci in range(n_cols):
                self._remove_cell_borders(table.cell(ri, ci))

        # Add teal accent line under header
        self._add_header_accent_line(table, n_cols)

        # Apply alternating fills on data rows
        if dark:
            for ri in range(1, n_rows):
                alpha = 8 if ri % 2 == 0 else 4
                for ci in range(n_cols):
                    cell = table.cell(ri, ci)
                    tc = cell._tc
                    tcPr = tc.find(qn('a:tcPr'))
                    if tcPr is None:
                        tcPr = etree.SubElement(tc, qn('a:tcPr'))
                    # Remove existing fill
                    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
                        existing = tcPr.find(qn(tag))
                        if existing is not None:
                            tcPr.remove(existing)
                    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                    srgb.set('val', 'FFFFFF')
                    alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                    alpha_elem.set('val', str(alpha * 1000))

        # CRITICAL: Reorder all tcPr children to match OOXML spec.
        # Without this, PowerPoint shows a repair dialog.
        for ri in range(n_rows):
            for ci in range(n_cols):
                cell = table.cell(ri, ci)
                tcPr = cell._tc.find(qn('a:tcPr'))
                if tcPr is not None:
                    self._reorder_tcPr_children(tcPr)

    def _style_dark_data_rows(self, table):
        """Style data rows for dark backgrounds.

        Alternating semi-transparent fills via OOXML XML.
        Status-colored text for numeric change columns.
        """
        for ri in range(1, len(table.rows)):
            alpha = 8 if ri % 2 == 0 else 5  # subtle alternation
            for ci in range(len(table.columns)):
                cell = table.cell(ri, ci)
                text = cell.text.strip()

                # Detect status-colored values
                text_color = COOL_GREY
                if _looks_numeric(text):
                    if text.startswith("-"):
                        text_color = HOT_RED
                    elif text.startswith("+"):
                        text_color = TEAL_GREEN

                # Apply semi-transparent fill via OOXML
                tc = cell._tc
                tcPr = tc.find(qn('a:tcPr'))
                if tcPr is None:
                    tcPr = etree.SubElement(tc, qn('a:tcPr'))
                # Remove existing fill
                for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
                    existing = tcPr.find(qn(tag))
                    if existing is not None:
                        tcPr.remove(existing)
                solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgb.set('val', 'FFFFFF')
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', str(alpha * 1000))

                # Style text
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = BRAND_FONT
                        run.font.size = Pt(14)
                        run.font.color.rgb = text_color
                    if not para.runs and para.text:
                        para.font.name = BRAND_FONT
                        para.font.size = Pt(14)
                        para.font.color.rgb = text_color

    @staticmethod
    def _auto_col_widths(headers, rows, total_width):
        """Column widths proportional to max content length (EMU)."""
        n = len(headers)
        if n == 0:
            return []
        lengths = [max(len(str(h)), 3) for h in headers]
        for row in rows:
            for ci, val in enumerate(row):
                if ci < n:
                    lengths[ci] = max(lengths[ci], len(str(val)))
        tw = int(total_width)
        total = sum(lengths)
        widths = [int(tw * l / total) for l in lengths]
        widths[-1] += tw - sum(widths)  # absorb rounding remainder
        return widths


_DEFAULT_TABLE_POS = {
    "left": Inches(1.0), "top": Inches(2.0),
    "width": Inches(11.0), "height": Inches(4.5),
}


def _looks_numeric(text: str) -> bool:
    """Heuristic: detect numbers like 123, $1.2M, 45.6%, +12%."""
    if not text:
        return False
    cleaned = text.replace(",", "").replace("$", "").replace("%", "")
    cleaned = cleaned.replace("+", "").replace("-", "").replace(" ", "").rstrip("KMBkmb")
    try:
        float(cleaned)
        return True
    except ValueError:
        return False


# ======================================================================
# DataFormatter
# ======================================================================
class DataFormatter:
    """Static utilities for number formatting in presentations.

    Usage: DataFormatter.format_currency(1234.5) -> "$1,235"
    """

    @staticmethod
    def format_currency(value, symbol="$", decimals=0):
        """1234.5 -> '$1,235' or '$1,234.50' with decimals=2."""
        fmt = f"{value:,.{decimals}f}"
        return f"{symbol}{fmt}"

    @staticmethod
    def format_percentage(value, decimals=1):
        """Auto-detect fraction vs percentage.
        0.153 -> '15.3%', 15.3 -> '15.3%', 1.0 -> '100.0%'
        """
        if value == 0:
            return f"0.{'0' * decimals}%"
        pct = value * 100 if -1 < value < 1 else value
        return f"{pct:.{decimals}f}%"

    @staticmethod
    def format_number(value, decimals=0):
        """1234567 -> '1,234,567'"""
        return f"{value:,.{decimals}f}"

    @staticmethod
    def format_large_number(value):
        """Abbreviate: 1234567 -> '1.2M', 1234 -> '1.2K', 500 -> '500'."""
        a = abs(value)
        sign = "-" if value < 0 else ""
        if a >= 1_000_000_000:
            return f"{sign}{a / 1e9:.1f}B"
        if a >= 1_000_000:
            return f"{sign}{a / 1e6:.1f}M"
        if a >= 1_000:
            return f"{sign}{a / 1e3:.1f}K"
        return f"{sign}{int(a)}" if a == int(a) else f"{sign}{a:.1f}"

    @staticmethod
    def format_change(value, positive_prefix="+"):
        """+5.2 -> '+5.2%', -3.1 -> '-3.1%', 0 -> '0.0%'."""
        if value > 0:
            return f"{positive_prefix}{value:.1f}%"
        if value < 0:
            return f"{value:.1f}%"
        return "0.0%"

    @staticmethod
    def prepare_table_data(raw_data, formatters=None):
        """Apply formatter functions to specific columns.

        Args:
            raw_data: List of row lists.
            formatters: {col_index: callable} mapping.
        Returns: List of row lists with all values as strings.

        Example:
            fmt = {1: DataFormatter.format_currency, 2: DataFormatter.format_percentage}
            DataFormatter.prepare_table_data(data, fmt)
        """
        formatters = formatters or {}
        result = []
        for row in raw_data:
            out = []
            for ci, val in enumerate(row):
                if ci in formatters:
                    try:
                        out.append(formatters[ci](val))
                    except (TypeError, ValueError):
                        out.append(str(val))
                else:
                    out.append(str(val))
            result.append(out)
        return result
