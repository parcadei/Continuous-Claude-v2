"""
PPTX Content Extractor
Usage: python pptx-extractor.py <input.pptx> <output_dir>

Extracts all slides, text, images, and speaker notes from a PowerPoint file.
Outputs a JSON structure suitable for HTML conversion.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os
import sys


def extract_pptx(file_path, output_dir):
    """
    Extract all content from a PowerPoint file.
    Returns a list of slide dicts with title, content, images, and notes.
    """
    prs = Presentation(file_path)
    slides_data = []

    assets_dir = os.path.join(output_dir, 'assets')
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            'number': slide_num + 1,
            'title': '',
            'content': [],
            'images': [],
            'notes': ''
        }

        for shape in slide.shapes:
            # Extract title
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data['title'] = shape.text
                else:
                    slide_data['content'].append({
                        'type': 'text',
                        'content': shape.text
                    })

            # Extract images (shape_type 13 = Picture)
            if shape.shape_type == 13:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                image_path = os.path.join(assets_dir, image_name)

                with open(image_path, 'wb') as f:
                    f.write(image_bytes)

                slide_data['images'].append({
                    'path': f"assets/{image_name}",
                    'width': shape.width,
                    'height': shape.height
                })

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data['notes'] = notes_frame.text

        slides_data.append(slide_data)

    return slides_data


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python pptx-extractor.py <input.pptx> <output_dir>")
        sys.exit(1)

    file_path = sys.argv[1]
    output_dir = sys.argv[2]
    os.makedirs(output_dir, exist_ok=True)

    slides = extract_pptx(file_path, output_dir)

    output_json = os.path.join(output_dir, 'slides.json')
    with open(output_json, 'w') as f:
        json.dump(slides, f, indent=2)

    print(f"Extracted {len(slides)} slides to {output_json}")
    for slide in slides:
        img_count = len(slide['images'])
        print(f"  Slide {slide['number']}: {slide['title']!r} ({img_count} image(s))")
